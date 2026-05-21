from __future__ import annotations

import re
import shutil
import sys
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]

OLD_EXPANSION_TITLES = {
    "代码组织与可复现运行",
    "误差分析与改进路线",
    "实现细节与关键权衡",
    "泛化表现与优化路线",
    "数据处理与训练闭环",
    "生成质量与采样策略",
    "讨论与展望",
}


@dataclass(frozen=True)
class Discussion:
    limitation: tuple[str, str]
    application: tuple[str, str]
    future: tuple[str, str]


@dataclass(frozen=True)
class ExperimentDeck:
    name: str
    title: str
    flow_title: str
    flow_subtitle: str
    flow_image: Path
    theme: tuple[int, int, int]
    accent: tuple[int, int, int]
    base: Path
    targets: tuple[Path, ...]
    discussion: Discussion
    summary_items: tuple[str, str, str]


EXPERIMENTS: tuple[ExperimentDeck, ...] = (
    ExperimentDeck(
        name="实验1",
        title="实验一：手写数字字符识别",
        flow_title="网络结构与训练流程",
        flow_subtitle="MNIST CNN · 特征提取、分类头与反向传播闭环",
        flow_image=ROOT / "实验报告/实验1/figures/algorithm_flow_20260426_vector.png",
        theme=(37, 99, 235),
        accent=(14, 165, 233),
        base=ROOT / "PPT/实验1/实验一PPT.pptx",
        targets=(
            ROOT / "PPT/实验1/实验一PPT.pptx",
            ROOT / "PPT/实验1/实验一PPT_svg.pptx",
            ROOT / "项目提交/实验1/实验PPT/实验1PPT.pptx",
            ROOT / "项目提交/实验1/实验PPT/实验1PPT_svg.pptx",
        ),
        discussion=Discussion(
            limitation=("局限性", "MNIST 背景干净、尺寸统一，无法覆盖真实票据、表格扫描和拍照场景中的噪声、倾斜、遮挡与分布漂移。"),
            application=("应用方向", "可扩展到答题卡识别、票据数字读取、仪表盘编号识别和轻量 OCR 系统中的局部字符分类模块。"),
            future=("深入研究", "后续应加入鲁棒数据增强、置信度校准、拒识机制、错例聚类和轻量化部署评估，而不仅追求更高测试准确率。"),
        ),
        summary_items=(
            "完成 PyTorch/CUDA 环境验证、CNN 模型训练、checkpoint 复评和可视化分析。",
            "测试准确率达到 99.28%，显著超过 98% 的课程验收目标。",
            "实验重点在于建立可复现的图像分类工程链路，而不仅是获得单次高分。",
        ),
    ),
    ExperimentDeck(
        name="实验2",
        title="实验二：利用ViT实现CIFAR10图像分类",
        flow_title="算法设计与训练流程",
        flow_subtitle="CIFAR-10 Vision Transformer · Patch tokenization 到分类概率",
        flow_image=ROOT / "实验报告/实验2/figures/algorithm_flow_academic.png",
        theme=(8, 126, 164),
        accent=(245, 158, 11),
        base=ROOT / "PPT/实验2/实验2 PPT.pptx",
        targets=(
            ROOT / "PPT/实验2/实验2 PPT.pptx",
            ROOT / "PPT/实验2/实验2 PPT_svg.pptx",
            ROOT / "项目提交/实验2/实验PPT/实验2PPT.pptx",
            ROOT / "项目提交/实验2/实验PPT/实验2PPT_svg.pptx",
        ),
        discussion=Discussion(
            limitation=("局限性", "从零训练 ViT 对数据规模、正则化和训练调度高度敏感；在小尺寸图像上，纯注意力结构并不天然具备 CNN 的局部归纳偏置。"),
            application=("应用方向", "适合进一步迁移到遥感图像、医学影像初筛、工业缺陷检测等需要全局上下文建模的视觉任务。"),
            future=("深入研究", "建议开展 patch 尺寸、层数、注意力头、标签平滑和增强策略的系统消融，并探索 DeiT 蒸馏、CNN-ViT 混合结构与 token pruning。"),
        ),
        summary_items=(
            "从零实现 Patch Embedding、Transformer Encoder、class token 与分类头。",
            "最终测试准确率达到 86.41%，超过 80% 的课程目标。",
            "ViT 在小图像上的效果依赖训练策略，后续应补充消融、鲁棒性和效率评估。",
        ),
    ),
    ExperimentDeck(
        name="实验3",
        title="实验3：自动写诗",
        flow_title="算法设计与生成流程",
        flow_subtitle="LSTM 自动写诗 · 语料编码、序列建模与温度采样",
        flow_image=ROOT / "实验报告/实验3/figures/algorithm_flow_academic.png",
        theme=(124, 58, 237),
        accent=(236, 72, 153),
        base=ROOT / "PPT/实验3/实验3自动写诗_20260504_193517.pptx",
        targets=(
            ROOT / "PPT/实验3/实验3自动写诗_20260504_193517.pptx",
            ROOT / "PPT/实验3/实验3自动写诗_20260504_193517_svg.pptx",
            ROOT / "项目提交/实验3/实验PPT/实验3PPT.pptx",
            ROOT / "项目提交/实验3/实验PPT/实验3PPT_svg.pptx",
        ),
        discussion=Discussion(
            limitation=("局限性", "字符级 LSTM 能学习局部字词转移，但对长距离语义、格律押韵、主题一致性和审美质量的控制仍然有限。"),
            application=("应用方向", "可用于古诗教学、创作灵感辅助、藏头诗生成、数字人文实验和传统文化内容生成原型。"),
            future=("深入研究", "后续可引入格律/韵脚约束、统计化人工评价、Transformer 解码器、检索增强和 human-in-the-loop 交互式改写。"),
        ),
        summary_items=(
            "完成唐诗语料处理、Embedding + 双层 LSTM 建模、训练和生成推理流程。",
            "验证损失收敛至 1.9350，并生成首句续写与藏头诗示例。",
            "模型具备基本风格学习能力，但文学质量仍需要约束解码与人工评价共同提升。",
        ),
    ),
    ExperimentDeck(
        name="实验4",
        title="实验四：采用Transformer架构的神经机器翻译系统",
        flow_title="算法设计与翻译流程",
        flow_subtitle="Transformer NMT · 编码解码、Beam Search 与 BLEU4 评估",
        flow_image=ROOT / "实验报告/实验4/figures/algorithm_flow_academic.png",
        theme=(37, 99, 235),
        accent=(245, 158, 11),
        base=ROOT / "PPT/实验4/实验4 PPT.pptx",
        targets=(
            ROOT / "PPT/实验4/实验4 PPT.pptx",
            ROOT / "PPT/实验4/实验4 PPT_svg.pptx",
            ROOT / "项目提交/实验4/实验PPT/实验4PPT.pptx",
            ROOT / "项目提交/实验4/实验PPT/实验4PPT_svg.pptx",
        ),
        discussion=Discussion(
            limitation=("局限性", "当前实验基于中等规模平行语料和词级词表，仍受未登录词、长句一致性、领域偏移和 BLEU 单一指标的限制。"),
            application=("应用方向", "可迁移到课程资料翻译、技术文档辅助翻译和垂直领域双语检索，为更完整的机器翻译系统提供可复现原型。"),
            future=("深入研究", "后续应引入 BPE/SentencePiece、Transformer 正则化消融、多指标人工评测和领域自适应，系统分析效率与质量权衡。"),
        ),
        summary_items=(
            "完成 Transformer 编码器-解码器、注意力、训练评估和翻译推理闭环。",
            "最终测试集 BLEU4 达到 25.84，显著超过课程目标 14。",
            "实验价值不只在达标，更在于形成可复现、可扩展的 NMT 工程基线。",
        ),
    ),
    ExperimentDeck(
        name="实验6",
        title="实验六：基于 SegNet 的街景分割",
        flow_title="算法设计与分割流程",
        flow_subtitle="SegNet 街景分割 · 数据准备、像素级评估与样例预测",
        flow_image=ROOT / "实验报告/实验6/figures/algorithm_flow_academic.png",
        theme=(8, 126, 164),
        accent=(34, 197, 94),
        base=ROOT / "PPT/实验6/实验6 PPT.pptx",
        targets=(
            ROOT / "PPT/实验6/实验6 PPT.pptx",
            ROOT / "PPT/实验6/实验6 PPT_svg.pptx",
            ROOT / "项目提交/实验6/实验PPT/实验6PPT.pptx",
            ROOT / "项目提交/实验6/实验PPT/实验6PPT_svg.pptx",
        ),
        discussion=Discussion(
            limitation=("局限性", "CamVid Tiny 规模较小且类别分布不均衡，轻量 SegNet 易偏向大面积区域，对小目标、边界和少数类别识别不足。"),
            application=("应用方向", "街景语义分割可服务自动驾驶、道路巡检、智能交通和城市空间理解，是场景结构化感知的重要基础模块。"),
            future=("深入研究", "建议扩展到完整数据集，加入类别重加权、强骨干、多尺度上下文和边界损失，并报告稳定性与实时性评估。"),
        ),
        summary_items=(
            "完成 CamVid 数据准备、SegNet 实现、训练评估、指标计算和预测可视化。",
            "真实输出像素准确率、平均像素准确率和平均交并比，形成可追溯实验证据。",
            "后续优化重点应从单次跑通转向类别不均衡、边界质量和跨场景泛化。",
        ),
    ),
    ExperimentDeck(
        name="实验7",
        title="实验七：神经网络语言模型",
        flow_title="算法设计与语言建模流程",
        flow_subtitle="PTB LSTM LM · BPTT、验证集选择与测试 PPL",
        flow_image=ROOT / "实验报告/实验7/figures/algorithm_flow.png",
        theme=(79, 70, 229),
        accent=(16, 185, 129),
        base=ROOT / "PPT/实验7/实验7 PPT.pptx",
        targets=(
            ROOT / "PPT/实验7/实验7 PPT.pptx",
            ROOT / "PPT/实验7/实验7 PPT_svg.pptx",
            ROOT / "项目提交/实验7/实验PPT/实验7PPT.pptx",
            ROOT / "项目提交/实验7/实验PPT/实验7PPT_svg.pptx",
        ),
        discussion=Discussion(
            limitation=("局限性", "词级 LSTM 依赖固定词表和截断 BPTT，难以充分建模开放词汇、超长上下文和复杂语义迁移。"),
            application=("应用方向", "语言模型能力可用于输入法候选排序、文本补全、语音识别后处理和轻量端侧序列建模。"),
            future=("深入研究", "后续应比较 LSTM、GRU 与 Transformer，系统评估 dropout、权重绑定、学习率调度和子词建模对 PPL 的影响。"),
        ),
        summary_items=(
            "完成 PTB 数据准备、词表构建、两层 LSTM 训练和独立测试评估。",
            "最终测试 PPL 为 75.77，低于课程要求阈值 80，实验结果达标。",
            "达标配置揭示了优化器、学习率衰减和梯度裁剪对语言模型训练稳定性的关键作用。",
        ),
    ),
)


def slide_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            parts.append(shape.text)
    return "\n".join(parts)


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    prs.slides._sldIdLst.remove(slide_id)


def move_slide_to_end(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    prs.slides._sldIdLst.remove(slide_id)
    prs.slides._sldIdLst.append(slide_id)


def move_slide_to_start(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    prs.slides._sldIdLst.remove(slide_id)
    prs.slides._sldIdLst.insert(0, slide_id)


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


def find_slide(prs: Presentation, keywords: tuple[str, ...]) -> int | None:
    for idx, slide in enumerate(prs.slides):
        text = slide_text(slide)
        if all(keyword in text for keyword in keywords):
            return idx
    return None


def find_slide_by_title(prs: Presentation, title: str) -> int | None:
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if shape.text.strip().splitlines()[0].strip() == title:
                    return idx
    return None


def find_summary_slide(prs: Presentation) -> int | None:
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text in {"总结", "总结与展望"}:
                    return idx
    return None


def normalize_cover_slide(prs: Presentation, exp: ExperimentDeck) -> None:
    for idx, slide in enumerate(prs.slides):
        text = slide_text(slide)
        is_cover = exp.title in text and "感谢聆听" not in text and ("课程" in text or "2026年04月" in text)
        if is_cover and idx != 0:
            move_slide_to_start(prs, idx)
            return


def remove_expansion_slides(prs: Presentation) -> None:
    for idx in range(len(prs.slides) - 1, -1, -1):
        text = slide_text(prs.slides[idx])
        is_standalone_thanks = "感谢聆听" in text and "总结" not in text and len(text) < 180
        if any(title in text for title in OLD_EXPANSION_TITLES) or is_standalone_thanks:
            delete_slide(prs, idx)


def add_textbox(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def add_background(slide, prs: Presentation, theme):
    sw, sh = prs.slide_width, prs.slide_height
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(248, 250, 252)
    bg.line.fill.background()
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, Inches(0.12))
    top.fill.solid()
    top.fill.fore_color.rgb = RGBColor(*theme)
    top.line.fill.background()


def add_title(slide, exp: ExperimentDeck, title: str, subtitle: str):
    add_textbox(slide, Inches(0.70), Inches(0.36), Inches(10.0), Inches(0.38), title, 24, (15, 23, 42), True)
    add_textbox(slide, Inches(0.70), Inches(0.80), Inches(10.8), Inches(0.26), subtitle, 10.5, (100, 116, 139))
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.52), Inches(0.41), Inches(0.08), Inches(0.55))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(*exp.accent)
    accent.line.fill.background()


def fit_picture(slide, image_path: Path, left, top, max_width, max_height):
    with Image.open(image_path) as img:
        aspect = img.width / img.height
    width = max_width
    height = int(width / aspect)
    if height > max_height:
        height = max_height
        width = int(height * aspect)
    slide.shapes.add_picture(str(image_path), left + int((max_width - width) / 2), top + int((max_height - height) / 2), width=width, height=height)


def refresh_flow_slide(prs: Presentation, exp: ExperimentDeck):
    idx = find_slide_by_title(prs, exp.flow_title)
    if idx is None:
        idx = min(3, len(prs.slides) - 1)
    slide = prs.slides[idx]
    clear_slide(slide)
    add_background(slide, prs, exp.theme)
    add_title(slide, exp, exp.flow_title, exp.flow_subtitle)
    fit_picture(slide, exp.flow_image, Inches(0.55), Inches(1.32), prs.slide_width - Inches(1.10), prs.slide_height - Inches(1.78))


def render_discussion_slide(slide, prs: Presentation, exp: ExperimentDeck):
    clear_slide(slide)
    add_background(slide, prs, exp.theme)
    add_title(slide, exp, "讨论与展望", "从实验结论走向局限分析、应用外延与后续研究")

    cards = [exp.discussion.limitation, exp.discussion.application, exp.discussion.future]
    lefts = [Inches(0.62), Inches(4.73), Inches(8.84)]
    fills = [(239, 246, 255), (240, 253, 244), (255, 247, 237)]
    for (heading, body), left, fill in zip(cards, lefts, fills):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.58), Inches(3.55), Inches(4.70))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*fill)
        card.line.color.rgb = RGBColor(203, 213, 225)
        card.line.width = Pt(1.0)
        add_textbox(slide, left + Inches(0.25), Inches(1.88), Inches(3.05), Inches(0.36), heading, 16, exp.theme, True, PP_ALIGN.CENTER)
        rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.58), Inches(2.38), Inches(2.40), Pt(1.4))
        rule.fill.solid()
        rule.fill.fore_color.rgb = RGBColor(*exp.accent)
        rule.line.fill.background()
        add_textbox(slide, left + Inches(0.33), Inches(2.68), Inches(2.90), Inches(2.80), body, 11.2, (51, 65, 85))

    add_textbox(
        slide,
        Inches(0.78),
        Inches(6.63),
        Inches(11.7),
        Inches(0.28),
        "核心判断：当前实验已经完成课程目标，但真正的研究价值来自对边界条件、应用迁移和后续验证路径的清醒说明。",
        10.3,
        (71, 85, 105),
        True,
        PP_ALIGN.CENTER,
    )


def render_thanks_slide(slide, prs: Presentation, exp: ExperimentDeck):
    clear_slide(slide)
    sw, sh = prs.slide_width, prs.slide_height
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(15, 23, 42)
    bg.line.fill.background()

    block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(1.18), Inches(0.10), Inches(1.25))
    block.fill.solid()
    block.fill.fore_color.rgb = RGBColor(*exp.accent)
    block.line.fill.background()
    add_textbox(slide, Inches(1.15), Inches(1.10), Inches(10.5), Inches(0.58), exp.title, 18, (226, 232, 240), True)
    add_textbox(slide, Inches(1.15), Inches(2.68), Inches(10.8), Inches(0.82), "感谢聆听！", 44, (255, 255, 255), True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.15), Inches(3.66), Inches(10.8), Inches(0.40), "敬请老师和同学批评指正", 18, (203, 213, 225), False, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.15), Inches(6.50), Inches(10.8), Inches(0.25), "Deep Learning Course Experiment Presentation", 9.5, (148, 163, 184), False, PP_ALIGN.CENTER)


def render_summary_slide(slide, prs: Presentation, exp: ExperimentDeck):
    clear_slide(slide)
    add_background(slide, prs, exp.theme)
    add_title(slide, exp, "总结", "实验目标、结果证据与核心收获")

    left = Inches(1.30)
    top = Inches(1.72)
    for idx, item in enumerate(exp.summary_items, 1):
        y = top + Inches((idx - 1) * 1.28)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, y, Inches(0.46), Inches(0.46))
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(*exp.accent)
        circle.line.fill.background()
        add_textbox(slide, left, y + Inches(0.06), Inches(0.46), Inches(0.22), str(idx), 11, (255, 255, 255), True, PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.72), y - Inches(0.03), Inches(9.70), Inches(0.58), item, 16, (30, 41, 59), idx == 1)


def normalize_summary_slide(prs: Presentation) -> bool:
    idx = find_summary_slide(prs)
    if idx is None:
        return False
    slide = prs.slides[idx]
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            text = shape.text.replace("总结与展望", "总结")
            text = text.replace("感谢聆听！", "").replace("谢谢聆听！", "")
            shape.text = text.strip()
    move_slide_to_end(prs, idx)
    return True


def update_page_numbers(prs: Presentation):
    total = len(prs.slides)
    pattern = re.compile(r"^\s*\d+\s*/\s*\d+\s*$")
    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and pattern.match(shape.text):
                shape.text = f"{idx} / {total}"


def ensure_slide_count(prs: Presentation, count: int) -> None:
    while len(prs.slides) > count:
        delete_slide(prs, len(prs.slides) - 1)
    while len(prs.slides) < count:
        prs.slides.add_slide(prs.slide_layouts[6])


def update_deck(source_path: Path, deck_path: Path, exp: ExperimentDeck) -> Path:
    prs = Presentation(BytesIO(source_path.read_bytes()))
    normalize_cover_slide(prs, exp)
    refresh_flow_slide(prs, exp)
    ensure_slide_count(prs, 10)
    render_discussion_slide(prs.slides[7], prs, exp)
    render_summary_slide(prs.slides[8], prs, exp)
    render_thanks_slide(prs.slides[9], prs, exp)
    update_page_numbers(prs)
    if len(prs.slides) != 10:
        raise RuntimeError(f"{deck_path} has {len(prs.slides)} slides after update; expected 10")
    tmp_path = deck_path.with_name(f".{deck_path.stem}.expanded.tmp.pptx")
    prs.save(str(tmp_path))
    try:
        shutil.copyfile(tmp_path, deck_path)
        written_path = deck_path
    except PermissionError:
        written_path = deck_path.with_name(f"{deck_path.stem}_10页优化版.pptx")
        shutil.copyfile(tmp_path, written_path)
        print(f"locked, wrote optimized copy instead: {written_path.relative_to(ROOT)}")
    finally:
        tmp_path.unlink(missing_ok=True)
    return written_path


def main():
    requested = set(sys.argv[1:]) if len(sys.argv) > 1 else {exp.name for exp in EXPERIMENTS}
    known = {exp.name for exp in EXPERIMENTS}
    unknown = requested.difference(known)
    if unknown:
        raise SystemExit(f"unknown experiment deck(s): {', '.join(sorted(unknown))}")

    updated: list[Path] = []
    for exp in EXPERIMENTS:
        if exp.name not in requested:
            continue
        if not exp.base.exists():
            print(f"skip {exp.name}: base deck missing: {exp.base.relative_to(ROOT)}")
            continue
        for target in exp.targets:
            if target.exists() or target.parent.exists():
                written = update_deck(exp.base, target, exp)
                updated.append(written)
                print(f"updated: {written.relative_to(ROOT)}")
            else:
                print(f"skip missing parent: {target.relative_to(ROOT)}")
    print(f"updated {len(updated)} deck(s)")


if __name__ == "__main__":
    main()
