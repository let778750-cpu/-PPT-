from __future__ import annotations

import argparse
import re
import shutil
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

from apply_latin_times_font import patch_pptx


EMU_PER_INCH = 914400
PAGE_RE = re.compile(r"^\s*\d{2}\s*/\s*\d{2}\s*$")
LATIN_RE = re.compile(r"[A-Za-z]")


@dataclass
class PolishResult:
    path: Path
    removed_shapes: int
    reflowed_text: int
    enlarged_text: int
    locked_output: Path | None = None


def inch_value(emu: int) -> float:
    return emu / EMU_PER_INCH


def has_text(shape) -> bool:
    return bool((getattr(shape, "text", "") or "").strip())


def shape_text(shape) -> str:
    return (getattr(shape, "text", "") or "").strip()


def is_page_number(shape) -> bool:
    text = shape_text(shape)
    return PAGE_RE.match(text) is not None and inch_value(shape.left) > 10.7 and inch_value(shape.top) > 6.4


def is_reference(shape) -> bool:
    text = shape_text(shape)
    return text.startswith("[") and inch_value(shape.left) < 1.0 and inch_value(shape.top) > 6.7


def should_remove_bottom_shape(slide_index: int, shape) -> bool:
    if slide_index == 1 or is_page_number(shape) or is_reference(shape):
        return False

    left = inch_value(shape.left)
    top = inch_value(shape.top)
    width = inch_value(shape.width)
    height = inch_value(shape.height)
    text = shape_text(shape)

    if top >= 6.0 and width >= 8.0 and text:
        return True

    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and not text:
        if top >= 6.25 and width >= 10.0 and height >= 0.28:
            return True
        if top >= 6.25 and width >= 10.0 and height <= 0.08:
            return True

    return False


def remove_shape(shape) -> None:
    element = shape._element  # noqa: SLF001 - python-pptx has no public delete API.
    element.getparent().remove(element)


def join_lines(text: str) -> str:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if len(lines) <= 1:
        return text

    joined = lines[0]
    for nxt in lines[1:]:
        prev = joined[-1] if joined else ""
        first = nxt[0]
        if LATIN_RE.match(prev) and LATIN_RE.match(first):
            sep = " "
        elif LATIN_RE.match(prev) and first not in "，。；、：,.!?;:)]}）】":
            sep = " "
        elif LATIN_RE.match(first) and prev not in "（([{":
            sep = " "
        else:
            sep = ""
        joined += sep + nxt
    return joined


def first_run_style(shape):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            return {
                "name": run.font.name,
                "size": run.font.size,
                "bold": run.font.bold,
                "italic": run.font.italic,
                "color": getattr(getattr(run.font, "color", None), "rgb", None),
            }
    return {"name": None, "size": None, "bold": None, "italic": None, "color": None}


def rewrite_text(shape, text: str, size_pt: float | None = None) -> None:
    style = first_run_style(shape)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = ""
    p._p.get_or_add_pPr().set("latinLnBrk", "0")  # noqa: SLF001
    run = p.add_run()
    run.text = text
    if style["name"]:
        run.font.name = style["name"]
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    elif style["size"]:
        run.font.size = style["size"]
    if style["bold"] is not None:
        run.font.bold = style["bold"]
    if style["italic"] is not None:
        run.font.italic = style["italic"]
    if style["color"]:
        run.font.color.rgb = style["color"]


def current_font_size_pt(shape) -> float | None:
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size:
                return run.font.size.pt
    return None


def target_body_size(shape, text: str) -> float | None:
    top = inch_value(shape.top)
    width = inch_value(shape.width)
    height = inch_value(shape.height)
    length = len(text)

    if top >= 3.0 and 2.4 <= width <= 3.4 and height >= 1.5:
        return 13.2 if length <= 45 else 12.6
    if top >= 2.2 and 2.2 <= width <= 3.0 and height >= 1.0:
        return 10.8 if length > 35 else 11.2
    if top >= 3.0 and width >= 4.0 and height >= 0.4:
        return 11.0
    return None


def should_reflow(shape) -> bool:
    if not hasattr(shape, "text_frame") or not has_text(shape):
        return False
    if is_page_number(shape) or is_reference(shape):
        return False

    text = shape_text(shape)
    top = inch_value(shape.top)
    width = inch_value(shape.width)
    height = inch_value(shape.height)

    if top < 2.0:
        return False
    if "•" in text:
        return False
    if len(text) < 12:
        return False
    return "\n" in text or (height >= 1.0 and width >= 2.4)


def polish_pptx(path: Path) -> PolishResult:
    prs = Presentation(path)
    removed = 0
    reflowed = 0
    enlarged = 0

    for slide_index, slide in enumerate(prs.slides, 1):
        for shape in list(slide.shapes):
            if should_remove_bottom_shape(slide_index, shape):
                remove_shape(shape)
                removed += 1

        for shape in list(slide.shapes):
            if not should_reflow(shape):
                continue
            original = shape_text(shape)
            new_text = join_lines(original)
            size = target_body_size(shape, new_text)
            old_size = current_font_size_pt(shape)
            if new_text != original or (size and (old_size is None or old_size < size - 0.05)):
                rewrite_text(shape, new_text, size)
                reflowed += int(new_text != original)
                enlarged += int(bool(size and (old_size is None or old_size < size - 0.05)))

    tmp = path.with_name(f".{path.stem}.polish.tmp.pptx")
    prs.save(tmp)
    patch_pptx(tmp)

    locked_output = None
    try:
        shutil.move(str(tmp), str(path))
    except PermissionError:
        locked_output = path.with_name(f"{path.stem}_polished_locked.pptx")
        shutil.move(str(tmp), str(locked_output))

    return PolishResult(path, removed, reflowed, enlarged, locked_output)


def main() -> None:
    parser = argparse.ArgumentParser(description="Remove empty bottom note frames and improve prose wrapping in PPTX files.")
    parser.add_argument("paths", nargs="+", type=Path)
    args = parser.parse_args()

    pptx_paths: list[Path] = []
    for path in args.paths:
        if path.is_dir():
            pptx_paths.extend(sorted(p for p in path.rglob("*.pptx") if not p.name.startswith("~$")))
        elif not path.name.startswith("~$"):
            pptx_paths.append(path)

    for path in pptx_paths:
        result = polish_pptx(path)
        suffix = f" | locked_output={result.locked_output}" if result.locked_output else ""
        print(
            f"polished | {result.path} | removed={result.removed_shapes} "
            f"| reflowed={result.reflowed_text} | enlarged={result.enlarged_text}{suffix}"
        )


if __name__ == "__main__":
    main()
