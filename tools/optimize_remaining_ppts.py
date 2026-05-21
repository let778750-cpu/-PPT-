from __future__ import annotations

import shutil
import re
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from optimize_experiment1_ppt import (
    BODY_FONT,
    CARD,
    CARD_2,
    GREEN,
    LIGHT_TAN,
    LINE,
    MONO_FONT,
    MUTED,
    ORANGE,
    SAGE,
    TAN,
    TEXT,
    TITLE_FONT,
    add_bg,
    add_card,
    add_footer_ref,
    add_header,
    add_inner,
    add_latex_formula,
    add_picture_fit,
    add_rect,
    add_text,
    prose_wrap,
)


ROOT = Path(__file__).resolve().parents[1]


@dataclass(frozen=True)
class FormulaCard:
    name: str
    title: str
    desc: str
    latex: str
    accent: object


@dataclass(frozen=True)
class ExperimentSpec:
    key: str
    title: str
    subtitle: str
    date: str
    outputs: tuple[Path, ...]
    overview_cards: tuple[tuple[str, str], ...]
    purpose: str
    formulas: tuple[FormulaCard, ...]
    theory_ref: str
    flow_title: str
    flow_subtitle: str
    flow_image: Path
    flow_note: str
    env_rows: tuple[tuple[str, str], ...]
    hp_rows: tuple[tuple[str, str], ...]
    setup_notes: tuple[tuple[str, str, object], ...]
    result_rows: tuple[tuple[str, str], ...]
    result_images: tuple[tuple[Path, str], ...]
    result_note: str
    analysis_metric: tuple[str, str, str, str]
    analysis_points: tuple[tuple[str, str], ...]
    discussion: tuple[tuple[str, str, object], ...]
    summary_items: tuple[tuple[str, str], ...]
    refs: tuple[str, str]


def p(rel: str) -> Path:
    return ROOT / rel


def soft_wrap(text: str, limit: int) -> str:
    """Insert stable manual line breaks without splitting Latin words."""
    wrapped: list[str] = []
    break_chars = "，。；、：,; "
    latin_token = re.compile(r"[A-Za-z][A-Za-z0-9_./+-]*")

    def safe_cut(line: str) -> int:
        boundary = min(len(line) - 1, limit)
        for match in latin_token.finditer(line):
            if match.start() < boundary < match.end():
                if match.end() <= limit + 8:
                    return match.end()
                if match.start() > 0:
                    return match.start()
                return match.end()
        for i in range(boundary, -1, -1):
            if line[i] in break_chars:
                return i + 1
        return boundary

    for raw in text.splitlines():
        line = raw.strip()
        while len(line) > limit:
            cut = safe_cut(line)
            if cut <= 0:
                cut = min(len(line), limit)
            wrapped.append(line[:cut].strip())
            line = line[cut:].strip()
        if line:
            wrapped.append(line)
    return "\n".join(wrapped)


EXPERIMENTS: tuple[ExperimentSpec, ...] = (
    ExperimentSpec(
        key="实验2",
        title="实验二：基于ViT的CIFAR10图像分类",
        subtitle="Vision Transformer · CIFAR-10 · 从零实现与达标验证",
        date="2026年4月",
        outputs=(p("PPT/实验2/实验2 PPT.pptx"), p("PPT/实验2/实验2 PPT_svg.pptx")),
        overview_cards=(
            ("实验目标", "从零实现 ViT 模型，\n在 CIFAR10 图像分类任务上\n达到 80% 以上测试准确率。"),
            ("数据集", "CIFAR-10\n50000 训练图像\n10000 测试图像\n32x32 彩色图像，10 类。"),
            ("模型架构", "Patch Embedding\nClass Token + Positional Embedding\n8 层 Transformer Encoder\nMLP 分类头。"),
            ("达标结果", "目标准确率 >= 80%\n最终测试准确率 86.41%\n超出目标 6.41 个百分点。"),
        ),
        purpose="本实验通过从零实现 ViT，验证 Transformer 在小尺寸图像分类中的建模能力；训练流程覆盖数据增强、正则化、学习率调度、checkpoint 选择与测试集复评。",
        formulas=(
            FormulaCard("vit_patch", "Patch 嵌入", "将图像切分为固定大小 patch，再映射为 token 序列并加入位置编码。", r"z_0=[x_{\mathrm{cls}};x_p^1E;\cdots;x_p^NE]+E_{\mathrm{pos}}", GREEN),
            FormulaCard("vit_attention", "缩放点积注意力", "通过查询、键和值的相似度分配权重，建模不同 patch 之间的全局依赖。", r"\mathrm{Attention}(Q,K,V)=\mathrm{softmax}\left(\frac{QK^\top}{\sqrt{d_k}}\right)V", SAGE),
            FormulaCard("vit_block", "Transformer Block", "多头注意力与 MLP 通过残差连接堆叠，逐层更新图像 token 表征。", r"\begin{aligned}z_\ell'&=\mathrm{MSA}(\mathrm{LN}(z_{\ell-1}))+z_{\ell-1}\\z_\ell&=\mathrm{MLP}(\mathrm{LN}(z_\ell'))+z_\ell'\end{aligned}", TAN),
            FormulaCard("vit_ce", "标签平滑损失", "在交叉熵中分配少量概率给非目标类，降低过度自信并改善泛化。", r"\mathcal{L}=-(1-\epsilon)\log p_y-\frac{\epsilon}{K}\sum_{k=1}^{K}\log p_k", GREEN),
        ),
        theory_ref="[1] Dosovitskiy et al., An Image is Worth 16x16 Words, ICLR 2021  [2] Vaswani et al., Attention Is All You Need, NeurIPS 2017",
        flow_title="算法设计与训练流程",
        flow_subtitle="Patch tokenization · Transformer 编码 · 分类头 · checkpoint 复评",
        flow_image=p("实验报告/实验2/figures/algorithm_flow_academic.png"),
        flow_note="输入图像经过数据增强后被切分为 patch 序列，经 Transformer Encoder 提取全局上下文，最终由 class token 完成 10 类分类。",
        env_rows=(("Python", "3.12.12"), ("PyTorch", "2.9.0+cu128"), ("CUDA", "12.8"), ("GPU", "RTX 5060 Laptop"), ("参数量", "6,350,602")),
        hp_rows=(("Epochs", "100"), ("Batch Size", "128"), ("LR", "3e-4"), ("Weight Decay", "0.05"), ("Warmup", "10"), ("Dropout", "0.1")),
        setup_notes=(("优化器", "AdamW + CosineAnnealingLR，配合 gradient clip=1.0 稳定训练。", GREEN), ("正则化", "RandomCrop、RandomHorizontalFlip、AutoAugment 与 Label Smoothing=0.1。", TAN)),
        result_rows=(("目标", "80%"), ("最佳测试准确率", "86.41%"), ("超出目标", "+6.41%"), ("训练时长", "约49分钟")),
        result_images=((p("实验报告/实验2/figures/training_overview.png"), "训练总览"), (p("实验报告/实验2/figures/accuracy_curve.png"), "准确率曲线")),
        result_note="最终测试准确率达到 86.41%，超过 80% 课程目标；最后若干 epoch 测试准确率在 86.2% 附近波动，说明模型收敛较稳定。",
        analysis_metric=("目标准确率", "80%", "实际准确率", "86.41%"),
        analysis_points=(("训练稳定性", "最后 5 个 epoch 保持在 86.2% 左右，未出现明显退化。"), ("关键贡献", "数据增强、标签平滑、warmup 与余弦退火共同提升泛化。"), ("研究边界", "小图像数据上 ViT 对训练策略更敏感，仍需消融验证。")),
        discussion=(
            ("局限性", "从零训练 ViT 对数据规模、正则化和训练调度高度敏感；小尺寸图像上纯注意力结构并不天然具备 CNN 的局部归纳偏置。", GREEN),
            ("应用方向", "可迁移到遥感图像、医学影像初筛、工业缺陷检测等需要全局上下文建模的视觉任务。", SAGE),
            ("深入研究", "后续应系统比较 patch 尺寸、层数、注意力头、标签平滑和增强策略，并探索 DeiT 蒸馏与 CNN-ViT 混合结构。", TAN),
        ),
        summary_items=(("1", "从零实现 Patch Embedding、Transformer Encoder、class token 与分类头。"), ("2", "最终测试准确率达到 86.41%，超过 80% 的课程目标。"), ("3", "实验说明 ViT 性能不仅来自结构，也高度依赖增强、正则化与调度策略。")),
        refs=("[3] Loshchilov & Hutter, Decoupled Weight Decay Regularization, ICLR 2019", "[4] CIFAR-10 dataset, Krizhevsky 2009"),
    ),
    ExperimentSpec(
        key="实验3",
        title="实验三：自动写诗",
        subtitle="LSTM 文本生成 · 唐诗语料 · 温度采样与藏头诗生成",
        date="2026年4月29日",
        outputs=(p("PPT/实验3/实验3自动写诗_20260504_193517.pptx"), p("PPT/实验3/实验3自动写诗_20260504_193517_svg.pptx")),
        overview_cards=(
            ("实验目标", "理解 RNN/LSTM 在深度学习框架中的实现，掌握文本生成流程。"),
            ("数据集", "tang.npz 唐诗语料\n57580 首诗\n词表规模 8293。"),
            ("模型结构", "Embedding + 2 层 LSTM\nDropout 正则化\n线性层输出字符概率。"),
            ("完成结果", "best val loss=1.9350\nfinal train loss=1.6690\n可生成续写诗和藏头诗。"),
        ),
        purpose="本实验使用字符级 LSTM 学习唐诗局部转移规律，并通过温度采样与 Top-k 控制生成多样性，重点验证文本生成工程链路而非虚构文学质量指标。",
        formulas=(
            FormulaCard("lstm_gates", "LSTM 门控", "遗忘门、输入门与输出门共同控制历史信息保留、新信息写入和当前输出。", r"\begin{aligned}f_t&=\sigma(W_f[h_{t-1},x_t]+b_f)\\i_t&=\sigma(W_i[h_{t-1},x_t]+b_i)\\o_t&=\sigma(W_o[h_{t-1},x_t]+b_o)\end{aligned}", GREEN),
            FormulaCard("lstm_cell", "记忆单元更新", "细胞状态以加性路径传递，有助于缓解长序列训练中的梯度衰减。", r"\begin{aligned}\tilde{c}_t&=\tanh(W_c[h_{t-1},x_t]+b_c)\\c_t&=f_t\odot c_{t-1}+i_t\odot \tilde{c}_t\end{aligned}", SAGE),
            FormulaCard("lm_loss", "序列交叉熵", "训练时根据历史字符预测下一个字符，最小化真实序列的负对数似然。", r"\mathcal{L}=-\frac{1}{T}\sum_{t=1}^{T}\log p_\theta(y_t\mid y_{<t})", TAN),
            FormulaCard("temperature", "温度采样", "温度参数调节概率分布尖锐程度，平衡生成的稳定性与多样性。", r"p_i=\frac{\exp(z_i/\tau)}{\sum_j\exp(z_j/\tau)}", GREEN),
        ),
        theory_ref="[1] Hochreiter & Schmidhuber, Long Short-Term Memory, Neural Computation 1997",
        flow_title="算法设计与生成流程",
        flow_subtitle="语料编码 · LSTM 建模 · 损失优化 · 续写与藏头诗生成",
        flow_image=p("实验报告/实验3/figures/algorithm_flow_academic.png"),
        flow_note="唐诗文本经词表编码后输入 2 层 LSTM，模型学习字符转移分布；推理阶段通过温度采样生成续写诗或藏头诗。",
        env_rows=(("Python", "3.12.12"), ("PyTorch", "2.9.0+cu128"), ("CUDA", "12.8"), ("GPU", "RTX 5060"), ("训练时间", "1499秒")),
        hp_rows=(("Embedding", "256"), ("Hidden", "512"), ("Layers", "2"), ("Dropout", "0.3"), ("Batch Size", "64"), ("LR", "0.001")),
        setup_notes=(("优化器", "Adam + StepLR，用于字符级语言模型训练。", GREEN), ("生成策略", "温度采样与 Top-k 控制候选空间，降低机械重复。", TAN)),
        result_rows=(("best val loss", "1.9350"), ("final train loss", "1.6690"), ("训练轮数", "30"), ("输出类型", "续写 / 藏头")),
        result_images=((p("实验报告/实验3/figures/loss_curve.png"), "损失曲线"), (p("实验报告/实验3/figures/model_architecture.png"), "模型结构")),
        result_note="验证损失收敛到 1.9350，模型能够生成基本符合汉语语法的诗句，但主题一致性和格律质量仍需要额外约束。",
        analysis_metric=("评价要求", "语法通顺", "生成结果", "通过"),
        analysis_points=(("续写能力", "模型能够围绕首句延续诗句结构，具备基本风格学习能力。"), ("藏头诗", "可按给定字头生成诗句，但语义连贯性仍不稳定。"), ("质量边界", "字符级 LSTM 更擅长局部转移，难以显式控制格律和主题。")),
        discussion=(
            ("局限性", "字符级 LSTM 能学习局部字词转移，但对长距离语义、格律押韵、主题一致性和审美质量的控制仍然有限。", GREEN),
            ("应用方向", "可用于古诗教学、创作灵感辅助、藏头诗生成、数字人文实验和传统文化内容生成原型。", SAGE),
            ("深入研究", "后续可引入格律/韵脚约束、统计化人工评价、Transformer 解码器和交互式改写。", TAN),
        ),
        summary_items=(("1", "完成唐诗语料处理、Embedding + 双层 LSTM 建模、训练和生成推理流程。"), ("2", "验证损失收敛至 1.9350，并生成首句续写与藏头诗示例。"), ("3", "模型具备基本风格学习能力，但文学质量仍需约束解码与人工评价共同提升。")),
        refs=("[2] PyTorch documentation and implementation evidence", "[3] Character-level language modeling practice"),
    ),
    ExperimentSpec(
        key="实验4",
        title="实验四：基于Transformer的神经机器翻译",
        subtitle="Transformer NMT · 中英平行语料 · Beam Search 与 BLEU4 评估",
        date="2026年4月30日",
        outputs=(p("PPT/实验4/实验4 PPT.pptx"), p("PPT/实验4/实验4 PPT_svg.pptx")),
        overview_cards=(
            ("实验目标", "使用 PyTorch 实现 Transformer 架构的中英机器翻译模型。"),
            ("数据集", "NiuTrans 中英平行语料\n100K 训练对\n400 dev 对\n1000 test 句。"),
            ("评价标准", "BLEU4 > 14\n使用 dev 集选择模型\n测试集独立评估。"),
            ("完成结果", "测试 BLEU4=25.84\n显著超过目标 14\n最佳 epoch=21。"),
        ),
        purpose="本实验从零实现 Transformer 编码器-解码器和翻译推理流程，重点验证注意力结构、标签平滑、warmup 调度与 Beam Search 在 NMT 中的作用。",
        formulas=(
            FormulaCard("nmt_attention", "缩放点积注意力", "注意力用源端与目标端表示的相似度加权聚合上下文，是 Transformer 的核心算子。", r"\mathrm{Attention}(Q,K,V)=\mathrm{softmax}\left(\frac{QK^\top}{\sqrt{d_k}}\right)V", GREEN),
            FormulaCard("multi_head", "多头注意力", "多个注意力头并行学习不同关系，再拼接映射为统一表示。", r"\mathrm{MultiHead}(Q,K,V)=\mathrm{Concat}(h_1,\ldots,h_h)W^O", SAGE),
            FormulaCard("position", "位置编码", "正弦位置编码为无循环结构注入词序位置信息。", r"\begin{aligned}\mathrm{PE}_{pos,2i}&=\sin(pos/10000^{2i/d})\\\mathrm{PE}_{pos,2i+1}&=\cos(pos/10000^{2i/d})\end{aligned}", TAN),
            FormulaCard("nmt_loss", "序列训练目标", "训练阶段最小化目标译文序列的平均负对数似然。", r"\mathcal{L}=-\frac{1}{T}\sum_{t=1}^{T}\log p_\theta(y_t\mid y_{<t},x)", GREEN),
        ),
        theory_ref="[1] Vaswani et al., Attention Is All You Need, NeurIPS 2017",
        flow_title="算法设计与翻译流程",
        flow_subtitle="编码器 · 解码器 · 标签平滑 · Beam Search · BLEU4 评估",
        flow_image=p("实验报告/实验4/figures/algorithm_flow_academic.png"),
        flow_note="源句经编码器生成上下文表示，解码器结合目标前缀逐词预测译文；推理阶段使用 Beam Search 改善全句搜索质量。",
        env_rows=(("Python", "3.12"), ("PyTorch", "2.9"), ("CUDA", "12.8"), ("GPU", "RTX 5060 8GB"), ("参数量", "17M")),
        hp_rows=(("Layers", "3+3"), ("d_model", "256"), ("Heads", "8"), ("Batch", "128"), ("LR", "5e-4"), ("Beam", "5")),
        setup_notes=(("优化策略", "Adam(β1=0.9, β2=0.98) + warmup 调度。", GREEN), ("训练目标", "标签平滑交叉熵降低过度自信，提高翻译泛化。", TAN)),
        result_rows=(("目标 BLEU4", ">14"), ("测试 BLEU4", "25.84"), ("最佳 epoch", "21"), ("验证损失", "1.9568")),
        result_images=((p("实验报告/实验4/figures/training_curves.png"), "训练曲线"), (p("实验报告/实验4/figures/principle_transformer_architecture.png"), "Transformer 原理")),
        result_note="验证集 BLEU4 从 17.43 提升到 25.88，测试集 BLEU4 为 25.84，说明模型稳定超过课程要求。",
        analysis_metric=("目标 BLEU4", ">14", "测试 BLEU4", "25.84"),
        analysis_points=(("关键设计", "Beam Search、标签平滑和 warmup 策略共同提升翻译质量。"), ("训练趋势", "验证 BLEU4 持续提升，后期进入稳定区间。"), ("改进空间", "词级词表仍受未登录词和领域偏移影响。")),
        discussion=(
            ("局限性", "当前实验基于中等规模平行语料和词级词表，仍受未登录词、长句一致性、领域偏移和 BLEU 单一指标限制。", GREEN),
            ("应用方向", "可迁移到课程资料翻译、技术文档辅助翻译和垂直领域双语检索。", SAGE),
            ("深入研究", "后续应引入 BPE/SentencePiece、多指标人工评测和领域自适应，系统分析效率与质量权衡。", TAN),
        ),
        summary_items=(("1", "完成 Transformer 编码器-解码器、注意力、训练评估和翻译推理闭环。"), ("2", "最终测试集 BLEU4 达到 25.84，显著超过课程目标 14。"), ("3", "实验价值在于形成可复现、可扩展的 NMT 工程基线。")),
        refs=("[2] Sennrich et al., Neural Machine Translation of Rare Words with Subword Units, 2016", "[3] Papineni et al., BLEU, ACL 2002"),
    ),
    ExperimentSpec(
        key="实验6",
        title="实验六：基于SegNet的街景分割",
        subtitle="CamVid Tiny · Encoder-Decoder · 像素级评估",
        date="2026年5月",
        outputs=(p("PPT/实验6/实验6 PPT.pptx"), p("PPT/实验6/实验6 PPT_svg.pptx")),
        overview_cards=(
            ("实验目标", "使用 Python 与 PyTorch\n构建街景语义分割模型。"),
            ("数据集", "CamVid Tiny 子集\n100 组图像/标签\n70/15/15 划分，32 类。"),
            ("评价指标", "像素准确率\n平均像素准确率\n平均交并比 mIoU。"),
            ("完成结果", "Pixel Acc=0.5362\nMean Acc=0.0977\nmIoU=0.0640。"),
        ),
        purpose="本实验实现紧凑 SegNet，覆盖语义分割数据准备、像素级训练、指标计算与样例预测，重点理解 Encoder-Decoder 和池化索引反池化机制。",
        formulas=(
            FormulaCard("seg_softmax", "像素级 Softmax", "对每个像素位置输出 32 类概率分布，用于逐像素分类。", r"p_{i,j,c}=\frac{\exp(z_{i,j,c})}{\sum_{k=1}^{C}\exp(z_{i,j,k})}", GREEN),
            FormulaCard("seg_ce", "分割交叉熵", "对所有有效像素的真实类别概率取负对数，形成训练目标。", r"\mathcal{L}=-\frac{1}{|\Omega|}\sum_{(i,j)\in\Omega}\log p_{i,j,y_{i,j}}", SAGE),
            FormulaCard("pixel_acc", "像素准确率", "统计预测类别与标签一致的像素比例，反映总体分类正确率。", r"\mathrm{PA}=\frac{\sum_c n_{cc}}{\sum_c t_c}", TAN),
            FormulaCard("miou", "平均交并比", "逐类别计算预测区域与真实区域的交并比，再对类别取平均。", r"\mathrm{mIoU}=\frac{1}{C}\sum_{c=1}^{C}\frac{n_{cc}}{t_c+\sum_k n_{kc}-n_{cc}}", GREEN),
        ),
        theory_ref="[1] Badrinarayanan et al., SegNet: A Deep Convolutional Encoder-Decoder Architecture, TPAMI 2017",
        flow_title="算法设计与分割流程",
        flow_subtitle="数据准备 · SegNet 前向 · 像素级损失 · 多指标评估",
        flow_image=p("实验报告/实验6/figures/algorithm_flow_academic.png"),
        flow_note="图像经编码器下采样提取语义特征，解码器利用反池化逐步恢复空间分辨率，最后输出每个像素的类别 logits。",
        env_rows=(("Python", "3.12.12"), ("PyTorch", "2.9.0+cu128"), ("CUDA", "12.8"), ("GPU", "RTX 5060 Laptop"), ("类别数", "32")),
        hp_rows=(("Epochs", "2"), ("Batch", "4"), ("LR", "1e-3"), ("Weight Decay", "1e-4"), ("输入尺寸", "128x96"), ("优化器", "AdamW")),
        setup_notes=(("损失函数", "像素级交叉熵，直接优化每个像素的类别预测。", GREEN), ("数据设置", "使用 CamVid Tiny 子集以保证课程环境可复现。", TAN)),
        result_rows=(("Pixel Acc", "0.5362"), ("Mean Acc", "0.0977"), ("mIoU", "0.0640"), ("报告要求", "已给出")),
        result_images=((p("实验报告/实验6/figures/training_curves.png"), "训练曲线"), (p("实验报告/实验6/figures/metrics_table.png"), "指标表")),
        result_note="模型已学习部分大面积结构，但小目标、边界和少数类别仍较弱；结果如实反映小数据、低轮数设置下的分割性能。",
        analysis_metric=("报告要求", "三项指标", "实际输出", "已完成"),
        analysis_points=(("主要表现", "大面积道路、天空等结构更容易被学习。"), ("薄弱区域", "小目标、边界和少数类别受数据规模与分辨率限制。"), ("解释原则", "指导书不设阈值，因此重点是指标完整和结果真实。")),
        discussion=(
            ("局限性", "CamVid Tiny 规模较小且类别分布不均衡，轻量 SegNet 易偏向大面积区域，对小目标和边界识别不足。", GREEN),
            ("应用方向", "街景语义分割可服务自动驾驶、道路巡检、智能交通和城市空间理解。", SAGE),
            ("深入研究", "后续应扩展完整数据集，加入类别重加权、强骨干、多尺度上下文和边界损失。", TAN),
        ),
        summary_items=(("1", "完成 CamVid 数据准备、SegNet 实现、训练评估、指标计算和预测可视化。"), ("2", "真实输出像素准确率、平均像素准确率和平均交并比，形成可追溯证据。"), ("3", "后续优化重点应转向类别不均衡、边界质量和跨场景泛化。")),
        refs=("[2] Brostow et al., CamVid dataset, 2008", "[3] Paszke et al., PyTorch, NeurIPS 2019"),
    ),
    ExperimentSpec(
        key="实验7",
        title="实验七：神经网络语言模型",
        subtitle="PTB LSTM LM · BPTT · 困惑度评估与达标配置",
        date="2026年5月15日",
        outputs=(p("PPT/实验7/实验7 PPT.pptx"), p("PPT/实验7/实验7 PPT_svg.pptx")),
        overview_cards=(
            ("实验目标", "构建标准 LSTM 网络，\n在 PTB 数据集上训练并评估语言模型。"),
            ("数据集", "PTB simple-examples\n词表 10000\ntrain/valid/test token 分别为 929589/73760/82430。"),
            ("评价指标", "困惑度 PPL\n越低表示模型对下一个词预测越确定。\n硬性目标：测试 PPL < 80。"),
            ("完成结果", "最佳验证 PPL=79.32\n测试 PPL=75.77\n最终结果达标。"),
        ),
        purpose="本实验通过 PTB LSTM 语言模型理解序列概率建模、截断 BPTT、困惑度评价以及优化器和学习率衰减对语言模型训练稳定性的影响。",
        formulas=(
            FormulaCard("lm_chain", "序列概率分解", "语言模型把整句概率分解为每个词在历史上下文条件下的概率乘积。", r"P(w_{1:T})=\prod_{t=1}^{T}P(w_t\mid w_{<t})", GREEN),
            FormulaCard("lm_nll", "负对数似然", "训练目标是最小化真实下一个词的平均负对数概率。", r"\mathcal{L}=-\frac{1}{T}\sum_{t=1}^{T}\log P_\theta(w_t\mid w_{<t})", SAGE),
            FormulaCard("ppl", "困惑度 PPL", "困惑度是平均损失的指数形式，数值越低表示预测越准确。", r"\mathrm{PPL}=\exp(\mathcal{L})", TAN),
            FormulaCard("lstm_update", "LSTM 状态更新", "门控记忆单元在长序列中保留必要历史信息，缓解梯度衰减。", r"\begin{aligned}c_t&=f_t\odot c_{t-1}+i_t\odot \tilde{c}_t\\h_t&=o_t\odot \tanh(c_t)\end{aligned}", GREEN),
        ),
        theory_ref="[1] Hochreiter & Schmidhuber, Long Short-Term Memory, Neural Computation 1997",
        flow_title="算法设计与语言建模流程",
        flow_subtitle="词表构建 · BPTT 训练 · 验证集选择 · 测试 PPL",
        flow_image=p("实验报告/实验7/figures/algorithm_flow.png"),
        flow_note="文本被编码为词 ID 序列后按 BPTT 长度展开训练，验证集选择最佳 checkpoint，最后在测试集报告 PPL。",
        env_rows=(("Python", "3.12.12"), ("PyTorch", "2.9.0+cu128"), ("CUDA", "12.8"), ("GPU", "RTX 5060 Laptop"), ("词表", "10000")),
        hp_rows=(("Epochs", "45"), ("Batch", "20"), ("Eval Batch", "10"), ("LR", "20"), ("Clip", "0.25"), ("BPTT", "35")),
        setup_notes=(("最终优化器", "SGD + 学习率衰减，是最终达标配置。", GREEN), ("调参过程", "初始 AdamW 25 轮测试 PPL=86.70，未达标后切换 SGD。", TAN)),
        result_rows=(("目标 PPL", "<80"), ("测试 PPL", "75.77"), ("最佳验证 PPL", "79.32"), ("最佳 epoch", "39")),
        result_images=((p("实验报告/实验7/figures/training_curves.svg"), "PPL 曲线"), (p("实验报告/实验7/figures/model_structure.png"), "模型结构")),
        result_note="最终测试 PPL 为 75.77，低于目标阈值 80；学习率衰减后验证 PPL 进一步下降，说明训练策略对达标至关重要。",
        analysis_metric=("目标 PPL", "<80", "测试 PPL", "75.77"),
        analysis_points=(("学习率衰减", "第 17 轮后学习率下降，验证 PPL 进一步改善。"), ("优化器选择", "SGD 更符合经典 PTB LSTM 配置，最终优于初始 AdamW。"), ("达标判断", "测试 PPL 75.77 < 80，满足课程要求。")),
        discussion=(
            ("局限性", "词级 LSTM 依赖固定词表和截断 BPTT，难以充分建模开放词汇、超长上下文和复杂语义迁移。", GREEN),
            ("应用方向", "语言模型能力可用于输入法候选排序、文本补全、语音识别后处理和轻量端侧序列建模。", SAGE),
            ("深入研究", "后续应比较 LSTM、GRU 与 Transformer，系统评估 dropout、权重绑定、学习率调度和子词建模对 PPL 的影响。", TAN),
        ),
        summary_items=(("1", "完成 PTB 数据准备、词表构建、两层 LSTM 训练和独立测试评估。"), ("2", "最终测试 PPL 为 75.77，低于课程阈值 80，实验结果达标。"), ("3", "达标配置揭示优化器、学习率衰减和梯度裁剪对语言模型训练稳定性的关键作用。")),
        refs=("[2] Zaremba et al., Recurrent Neural Network Regularization, ICLR 2015", "[3] Marcus et al., Penn Treebank, 1993"),
    ),
    ExperimentSpec(
        key="实验8",
        title="实验八：基于CNN-Transformer的图像描述",
        subtitle="MSCOCO 采样子集 · CNN 编码器 · Transformer 解码器",
        date="2026年5月18日",
        outputs=(p("PPT/实验8/实验8PPT.pptx"), p("PPT/实验8/实验8PPT_svg.pptx")),
        overview_cards=(
            ("实验目标", "构建 CNN 视觉编码器与 Transformer 文本解码器，完成图像到文本生成。"),
            ("数据集", "MSCOCO 2017 采样子集\n训练 360 图像/1800 描述\n验证/测试各 60 图像。"),
            ("模型结构", "CNN 提取视觉 token\nTransformer decoder 条件生成\nGreedy decoding。"),
            ("评价指标", "BLEU-1/4、METEOR、ROUGE-L、CIDEr\n指导书不设置硬性阈值。"),
        ),
        purpose="本实验使用真实 MSCOCO 图像描述数据而非合成样例，完整覆盖数据采样、CNN 预训练、条件语言建模、贪心解码和多指标评价。",
        formulas=(
            FormulaCard("caption_chain", "条件序列分解", "图像描述把整句概率分解为每一步在图像和历史词条件下的生成概率。", r"P(y_{1:T}\mid I)=\prod_{t=1}^{T}P(y_t\mid y_{<t},I)", GREEN),
            FormulaCard("caption_loss", "教师强制损失", "训练阶段使用真实前缀预测下一个词，最小化平均负对数似然。", r"\mathcal{L}(\theta)=-\frac{1}{T}\sum_{t=1}^{T}\log P_\theta(y_t\mid y_{<t},I)", SAGE),
            FormulaCard("caption_conv", "CNN 卷积特征", "卷积编码器从局部邻域提取视觉特征，为解码器提供图像条件。", r"Y_k(i,j)=\sigma\left(\sum_c\sum_{u,v}W_{k,c,u,v}X_c(i+u,j+v)+b_k\right)", TAN),
            FormulaCard("caption_attention", "交叉注意力", "Transformer 解码器通过注意力从视觉 token 中读取与当前词相关的图像信息。", r"\mathrm{Attention}(Q,K,V)=\mathrm{softmax}\left(\frac{QK^\top}{\sqrt{d_k}}\right)V", GREEN),
        ),
        theory_ref="[1] Vinyals et al., Show and Tell, CVPR 2015  [2] Vaswani et al., Attention Is All You Need, NeurIPS 2017",
        flow_title="算法设计与图像描述流程",
        flow_subtitle="MSCOCO 采样 · CNN 预训练 · Transformer 解码 · 多指标评估",
        flow_image=p("实验报告/实验8/figures/algorithm_flow.png"),
        flow_note="图像经 CNN 编码为视觉 token，文本前缀输入 Transformer decoder，模型逐词生成描述并用多参考指标评价。",
        env_rows=(("Python", "3.12.12"), ("PyTorch", "2.10.0+cpu"), ("CUDA", "不可用"), ("图像尺寸", "64x64"), ("词表", "1951")),
        hp_rows=(("Train images", "360"), ("Val/Test", "60/60"), ("Max Len", "24"), ("Decoder", "2 layers"), ("Heads", "4"), ("LR", "8e-4")),
        setup_notes=(("预训练", "基于 COCO instances 主目标类别对 CNN 编码器进行弱监督初始化。", GREEN), ("解码方式", "使用 greedy decoding，强调可复现和流程完整。", TAN)),
        result_rows=(("BLEU-1", "0.3261"), ("BLEU-4", "0.0070"), ("METEOR", "0.3515"), ("ROUGE-L", "0.3599"), ("CIDEr", "0.1714")),
        result_images=((p("实验报告/实验8/figures/training_curves.png"), "训练曲线"), (p("实验报告/实验8/figures/metric_table.png"), "指标表")),
        result_note="模型能学习部分常见词汇和语言模板，但在真实 MSCOCO 开放场景中视觉-语言对齐仍有限，BLEU-4 和 CIDEr 较低。",
        analysis_metric=("评价要求", "报告指标", "实际输出", "五项完整"),
        analysis_points=(("定量结果", "BLEU-1、METEOR、ROUGE-L 存在一定词汇重合，BLEU-4 与 CIDEr 暴露长程组合不足。"), ("定性样例", "预测容易收缩到高频模板，说明语言先验强于图像细节读取。"), ("实验边界", "CPU 小模型和 64x64 输入适合复现，但限制开放域描述质量。")),
        discussion=(
            ("局限性", "采样数据规模较小、图像分辨率低、CNN 弱监督预训练目标粗糙，贪心解码也会放大高频语言模板。", GREEN),
            ("应用方向", "图像描述可服务无障碍辅助、图像检索、教育资源标注、遥感理解和工业巡检摘要。", SAGE),
            ("深入研究", "后续应扩大数据规模，采用 ResNet/Swin/CLIP 等预训练视觉编码器，引入 beam search 与官方 COCO-caption 评价。", TAN),
        ),
        summary_items=(("1", "完成 MSCOCO 真实数据采样、CNN 预训练、Transformer 解码器训练和多指标评价。"), ("2", "BLEU-1=0.3261、METEOR=0.3515、ROUGE-L=0.3599，指标完整可追溯。"), ("3", "实验真实展示了图像描述从跑通流程到高质量开放域生成之间的难度差距。")),
        refs=("[3] Lin et al., Microsoft COCO: Common Objects in Context, ECCV 2014", "[4] Paszke et al., PyTorch, NeurIPS 2019"),
    ),
)


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def blank_layout(prs: Presentation):
    return prs.slide_layouts[6]


def page(num: int, total: int) -> str:
    return f"{num:02d} / {total:02d}"


def add_metric_pill(slide, label: str, value: str, x: float, y: float, w: float, accent=GREEN):
    add_rect(slide, x, y, w, 0.72, CARD_2, LINE, radius=True)
    add_text(slide, label, x + 0.16, y + 0.12, w - 0.32, 0.18, 8.6, MUTED, align=PP_ALIGN.CENTER, margin=0)
    add_text(slide, value, x + 0.12, y + 0.34, w - 0.24, 0.26, 15.2, accent, font=TITLE_FONT, bold=True, align=PP_ALIGN.CENTER, margin=0)


def add_table_rows(slide, rows: tuple[tuple[str, str], ...], x: float, y: float, w: float, row_h: float, value_w: float = 1.25):
    for i, (k, v) in enumerate(rows):
        row_y = y + i * row_h
        add_text(slide, k, x, row_y, w - value_w - 0.15, 0.22, 9.8, MUTED, font=BODY_FONT)
        add_text(slide, v, x + w - value_w, row_y, value_w, 0.23, 10.5, GREEN, bold=True, align=PP_ALIGN.RIGHT)
        add_rect(slide, x, row_y + 0.33, w, 0.006, LINE, None, radius=False)


def build_cover(prs: Presentation, spec: ExperimentSpec, total: int):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide, "cover")
    add_text(slide, "DEEP LEARNING COURSE", 0.47, 0.72, 3.6, 0.24, 9.2, MUTED, font=MONO_FONT)
    add_rect(slide, 0.47, 1.06, 1.05, 0.03, GREEN, None, radius=False)
    add_text(slide, spec.title, 0.42, 1.62, 8.55, 0.78, 28.0, GREEN, font=TITLE_FONT)
    add_text(slide, spec.subtitle, 0.47, 2.56, 7.2, 0.32, 14.2, TEXT)
    add_rect(slide, 0.47, 3.30, 5.8, 0.68, CARD, LINE, radius=True)
    add_text(slide, f"深度学习课程实验汇报 · {spec.date}", 0.74, 3.54, 5.25, 0.22, 10.8, MUTED, font=BODY_FONT, align=PP_ALIGN.CENTER)
    add_text(slide, spec.key.upper(), 0.48, 6.15, 1.4, 0.24, 9.5, MUTED, font=MONO_FONT)
    add_text(slide, "模型实现 · 训练评估 · 结果分析 · 总结展望", 2.0, 6.15, 4.2, 0.24, 9.5, MUTED)
    add_text(slide, page(1, total), 11.55, 6.78, 1.1, 0.25, 8.8, MUTED, font=MONO_FONT, align=PP_ALIGN.RIGHT)


def build_overview(prs: Presentation, spec: ExperimentSpec, total: int):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_header(slide, "EXPERIMENT OVERVIEW", "实验概述", "目标要求 · 数据集 · 模型配置 · 完成情况")
    accents = (GREEN, SAGE, TAN, GREEN)
    for i, (title, body) in enumerate(spec.overview_cards):
        x = 0.42 + i * 3.16
        add_card(slide, x, 1.78, 2.9, 2.65, title, accents[i])
        add_text(slide, f"{i+1:02d}", x + 2.22, 2.02, 0.5, 0.34, 14, accents[i], font=MONO_FONT, align=PP_ALIGN.RIGHT)
        add_text(slide, prose_wrap(body, 14.5), x + 0.18, 2.33, 2.56, 1.70, 10.2, TEXT, line_spacing=1.04)
    add_card(slide, 0.42, 4.72, 12.47, 1.82, "实验目的与汇报范围", GREEN, fill=CARD_2, pale=True)
    add_text(slide, soft_wrap(spec.purpose, 44), 0.68, 5.18, 11.9, 1.06, 12.0, TEXT, line_spacing=1.06)
    add_footer_ref(slide, spec.refs[1], page(2, total))


def build_theory(prs: Presentation, spec: ExperimentSpec, total: int):
    cards = spec.formulas
    chunks = [cards[i : i + 4] for i in range(0, len(cards), 4)]
    for chunk_index, chunk in enumerate(chunks, 1):
        slide = prs.slides.add_slide(blank_layout(prs))
        add_bg(slide)
        suffix = f"（{chunk_index}）" if len(chunks) > 1 else ""
        add_header(slide, "THEORETICAL BASIS", f"理论基础{suffix}", "核心公式 · 科学解释 · 与实验实现对应")
        positions = ((0.55, 1.76), (6.75, 1.76), (0.55, 4.18), (6.75, 4.18))
        for card, (x, y) in zip(chunk, positions):
            add_rect(slide, x, y, 5.75, 2.14, CARD, LINE, radius=True)
            add_rect(slide, x, y, 5.75, 0.045, card.accent, card.accent, radius=False)
            add_text(slide, card.title, x + 0.28, y + 0.20, 2.6, 0.28, 16.2, card.accent, font=TITLE_FONT, bold=True)
            add_text(slide, soft_wrap(card.desc, 24), x + 0.28, y + 0.60, 5.18, 0.56, 10.3, MUTED, line_spacing=1.0)
            add_latex_formula(slide, card.latex, x + 0.34, y + 1.24, 5.07, 0.76, card.accent, f"{spec.key}_{card.name}")
        add_footer_ref(slide, spec.theory_ref, page(3 + chunk_index - 1, total))


def build_flow(prs: Presentation, spec: ExperimentSpec, slide_no: int, total: int):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_header(slide, "MODEL PIPELINE", spec.flow_title, spec.flow_subtitle, page(slide_no, total))
    add_card(slide, 0.55, 1.76, 12.25, 4.85, None, GREEN)
    add_picture_fit(slide, spec.flow_image, 0.72, 1.96, 11.92, 4.46, pad=0.02, trim=True)


def build_setup(prs: Presentation, spec: ExperimentSpec, slide_no: int, total: int):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_header(slide, "EXPERIMENT SETUP", "实验环境与超参数配置", "软硬件环境 · 训练参数 · 优化目标")
    add_card(slide, 0.55, 1.78, 5.85, 3.85, "运行环境", GREEN)
    add_table_rows(slide, spec.env_rows, 0.9, 2.42, 4.9, 0.48, value_w=1.85)
    add_card(slide, 6.75, 1.78, 5.75, 3.85, "训练超参数", SAGE)
    add_table_rows(slide, spec.hp_rows, 7.08, 2.42, 4.55, 0.48, value_w=1.55)
    for i, (title, body, accent) in enumerate(spec.setup_notes[:2]):
        x = 0.55 + i * 6.2
        add_inner(slide, x, 5.82, 5.85 if i == 0 else 5.75, 0.86, body, 8.8, title, accent)
    add_footer_ref(slide, spec.refs[0], page(slide_no, total))


def build_results(prs: Presentation, spec: ExperimentSpec, slide_no: int, total: int):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_header(slide, "TRAINING RESULTS", "训练结果与验证数据", "关键指标 · 曲线证据 · 达标状态", page(slide_no, total))
    add_card(slide, 0.55, 1.78, 4.05, 4.82, "关键指标", GREEN)
    rows = spec.result_rows[:5]
    for i, (label, value) in enumerate(rows):
        add_metric_pill(slide, label, value, 0.88, 2.36 + i * 0.72, 3.40, GREEN if i in (1, 2) else SAGE)
    img1, label1 = spec.result_images[0]
    add_card(slide, 4.92, 1.78, 7.58, 4.82, label1, SAGE)
    add_picture_fit(slide, img1, 5.12, 2.18, 7.18, 3.92, pad=0.03, trim=True)


def build_evidence(prs: Presentation, spec: ExperimentSpec, image_item: tuple[Path, str], slide_no: int, total: int):
    image_path, label = image_item
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_header(slide, "SUPPORTING FIGURE", "补充图表与结果证据", "模型结构 · 指标表 · 样例预测 · 曲线补充", page(slide_no, total))
    add_card(slide, 0.55, 1.78, 11.95, 5.12, label, TAN)
    add_picture_fit(slide, image_path, 0.78, 2.38, 11.48, 4.20, pad=0.03, trim=True)


def build_analysis(prs: Presentation, spec: ExperimentSpec, slide_no: int, total: int):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_header(slide, "RESULT ANALYSIS", "结果分析与达标判定", "目标对比 · 关键发现 · 误差边界", page(slide_no, total))
    add_card(slide, 0.55, 1.78, 5.55, 4.55, "目标与实际", GREEN)
    l1, v1, l2, v2 = spec.analysis_metric
    add_text(slide, l1, 1.15, 2.33, 1.6, 0.24, 10.0, MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, v1, 1.05, 2.68, 1.8, 0.62, 25, MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "→", 3.02, 2.76, 0.5, 0.42, 17, GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, l2, 3.83, 2.33, 1.75, 0.24, 10.0, MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, v2, 3.72, 2.68, 1.95, 0.62, 25, GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.04, 3.65, 4.6, 0.006, LINE, None, radius=False)
    add_text(slide, "判定结论", 1.1, 3.86, 1.6, 0.28, 15, GREEN, font=TITLE_FONT, bold=True)
    add_text(slide, "实验结果与课程要求一致，\n关键指标已经完成可追溯验证。", 1.1, 4.20, 4.35, 0.52, 11.2, TEXT)
    add_rect(slide, 1.1, 4.95, 4.25, 0.50, LIGHT_TAN, None, radius=True)
    add_text(slide, "PASS", 1.1, 5.07, 4.25, 0.22, 11.5, GREEN, font=MONO_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, 6.35, 1.78, 6.15, 4.55, "关键发现", SAGE)
    for i, (title, body) in enumerate(spec.analysis_points):
        y = 2.36 + i * 1.17
        add_rect(slide, 6.78, y, 0.48, 0.48, CARD_2, LINE, radius=True)
        add_text(slide, f"{i+1}", 6.78, y + 0.12, 0.48, 0.18, 9.4, GREEN, font=MONO_FONT, bold=True, align=PP_ALIGN.CENTER, margin=0)
        add_text(slide, title, 7.45, y - 0.02, 1.8, 0.26, 13.6, TAN if i == 1 else GREEN, font=TITLE_FONT, bold=True)
        add_text(slide, soft_wrap(body, 25), 7.45, y + 0.30, 4.55, 0.48, 9.4, TEXT, line_spacing=0.95)


def build_discussion(prs: Presentation, spec: ExperimentSpec, slide_no: int, total: int):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_header(slide, "DISCUSSION & OUTLOOK", "讨论与展望", "从实验结论走向局限分析、应用外延与后续研究", page(slide_no, total))
    for i, (title, body, accent) in enumerate(spec.discussion):
        x = 0.45 + i * 4.25
        add_card(slide, x, 1.78, 3.95, 5.18, None, accent)
        add_rect(slide, x + 0.28, 2.1, 0.66, 0.66, CARD_2, LINE, radius=True)
        add_text(slide, f"{i+1}", x + 0.28, 2.24, 0.66, 0.32, 14.0, accent, font=MONO_FONT, bold=True, align=PP_ALIGN.CENTER, margin=0)
        add_text(slide, title, x + 1.08, 2.16, 2.25, 0.36, 18.0, accent, font=TITLE_FONT, bold=True)
        add_inner(slide, x + 0.09, 3.02, 3.77, 2.92, prose_wrap(body, 20.0), 12.1, accent=accent, fill=CARD_2, text_margin=0.08)


def build_summary(prs: Presentation, spec: ExperimentSpec, slide_no: int, total: int):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_header(slide, "SUMMARY", "总结", "实验目标、结果证据与核心收获", page(slide_no, total))
    for i, (num, body) in enumerate(spec.summary_items):
        x = 0.45 + i * 4.25
        accent = (GREEN, SAGE, TAN)[i]
        add_card(slide, x, 1.78, 3.95, 5.18, None, accent)
        add_rect(slide, x + 0.28, 2.1, 0.66, 0.66, CARD_2, LINE, radius=True)
        add_text(slide, num, x + 0.28, 2.24, 0.66, 0.32, 14.0, accent, font=MONO_FONT, bold=True, align=PP_ALIGN.CENTER, margin=0)
        add_text(slide, ["工程闭环", "结果证据", "核心收获"][i], x + 1.08, 2.16, 2.0, 0.36, 18.0, accent, font=TITLE_FONT, bold=True)
        add_inner(slide, x + 0.09, 3.02, 3.77, 2.92, prose_wrap(body, 20.0), 12.6, accent=accent, fill=CARD_2, text_margin=0.08)


def build_closing(prs: Presentation, spec: ExperimentSpec, total: int):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide, "close")
    add_rect(slide, 6.23, 0.62, 0.86, 0.86, CARD_2, LINE, radius=True)
    add_text(slide, spec.key.upper(), 6.25, 0.88, 0.82, 0.32, 10.2, GREEN, font=MONO_FONT, bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_text(slide, spec.title, 2.15, 1.95, 9.0, 0.65, 29, GREEN, font=TITLE_FONT, align=PP_ALIGN.CENTER)
    add_text(slide, "感谢聆听！", 4.1, 2.74, 5.2, 0.78, 30, GREEN, font=TITLE_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 3.75, 3.72, 1.35, 0.02, GREEN, None, radius=False)
    add_text(slide, "敬请老师和同学批评指正", 5.3, 3.56, 2.8, 0.36, 15, TEXT, font=TITLE_FONT, align=PP_ALIGN.CENTER)
    add_rect(slide, 8.25, 3.72, 1.35, 0.02, GREEN, None, radius=False)
    add_rect(slide, 3.65, 4.6, 6.1, 0.74, CARD, LINE, radius=True)
    add_text(slide, "Deep Learning Course Experiment Presentation", 3.9, 4.86, 5.6, 0.23, 11, MUTED, font=MONO_FONT, align=PP_ALIGN.CENTER)
    add_text(slide, page(total, total), 11.55, 6.82, 1.1, 0.24, 8.8, MUTED, font=MONO_FONT, align=PP_ALIGN.RIGHT)


def build_deck(spec: ExperimentSpec) -> Presentation:
    formula_slides = (len(spec.formulas) + 3) // 4
    evidence_slides = max(0, len(spec.result_images) - 1)
    total = 9 + formula_slides + evidence_slides
    prs = new_deck()
    build_cover(prs, spec, total)
    build_overview(prs, spec, total)
    build_theory(prs, spec, total)
    slide_no = 3 + formula_slides
    build_flow(prs, spec, slide_no, total)
    build_setup(prs, spec, slide_no + 1, total)
    result_no = slide_no + 2
    build_results(prs, spec, result_no, total)
    next_no = result_no + 1
    for image_item in spec.result_images[1:]:
        build_evidence(prs, spec, image_item, next_no, total)
        next_no += 1
    build_analysis(prs, spec, next_no, total)
    build_discussion(prs, spec, next_no + 1, total)
    build_summary(prs, spec, next_no + 2, total)
    build_closing(prs, spec, total)
    return prs


def save_to_outputs(spec: ExperimentSpec) -> list[Path]:
    for formula in spec.formulas:
        if not formula.latex.strip():
            raise ValueError(f"{spec.key}: empty formula {formula.title}")
    for img in [spec.flow_image, *(path for path, _ in spec.result_images)]:
        if not img.exists():
            raise FileNotFoundError(f"{spec.key}: missing image {img}")

    prs = build_deck(spec)
    main = spec.outputs[0]
    main.parent.mkdir(parents=True, exist_ok=True)
    tmp = main.with_name(f".{main.stem}.optimized.tmp.pptx")
    prs.save(tmp)
    written: list[Path] = []
    try:
        shutil.copyfile(tmp, main)
        written.append(main)
    except PermissionError:
        fallback = main.with_name(f"{main.stem}_优化版.pptx")
        shutil.copyfile(tmp, fallback)
        written.append(fallback)
    for out in spec.outputs[1:]:
        out.parent.mkdir(parents=True, exist_ok=True)
        try:
            shutil.copyfile(written[0], out)
            written.append(out)
        except PermissionError:
            fallback = out.with_name(f"{out.stem}_优化版.pptx")
            shutil.copyfile(written[0], fallback)
            written.append(fallback)
    tmp.unlink(missing_ok=True)
    return written


def main():
    for spec in EXPERIMENTS:
        outputs = save_to_outputs(spec)
        for out in outputs:
            print(f"updated {spec.key}: {out.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
