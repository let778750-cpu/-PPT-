from __future__ import annotations

import io
import re
import subprocess
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt
from reportlab.graphics import renderPM
from svglib.svglib import svg2rlg


ROOT = Path(__file__).resolve().parents[1]
OUT_ASSETS = ROOT / "_tmp_preview" / "optimized_ppt_assets"

BG = RGBColor(248, 247, 241)
CARD = RGBColor(255, 255, 252)
CARD_2 = RGBColor(241, 244, 238)
GREEN = RGBColor(38, 91, 61)
FORMULA_GREEN = RGBColor(18, 88, 55)
SAGE = RGBColor(126, 154, 124)
TAN = RGBColor(151, 124, 84)
LIGHT_TAN = RGBColor(245, 240, 232)
TEXT = RGBColor(52, 51, 46)
MUTED = RGBColor(119, 119, 111)
LINE = RGBColor(214, 218, 208)
ORANGE = RGBColor(193, 116, 44)

TITLE_FONT = "SimSun"
BODY_FONT = "Microsoft YaHei"
MONO_FONT = "Cascadia Mono"
LATIN_FONT = "Times New Roman"
LATIN_RE = re.compile(r"[A-Za-z]")
NO_LINE_START = set("\uFF0C\u3002\uFF1B\uFF1A\u3001,.!?)]}\uFF09\u3011\u300B%")
NO_LINE_END = set("([{\uFF08\u3010\u300A")
PROSE_BREAKS = set("\uFF0C\u3002\uFF1B\uFF1A\u3001,; \u7684\u548C\u4E0E\u53CA\u5E76")


def figure_path(name: str) -> Path:
    """Resolve experiment-1 report figures by filename, independent of PPT image order."""
    matches = list(ROOT.rglob(name))
    if not matches:
        raise FileNotFoundError(f"Experiment 1 figure not found: {name}")
    return matches[0]


def locate_source() -> tuple[Path, Path]:
    sample = next(ROOT.glob("AI *.pptx"))
    source = [
        p
        for p in ROOT.glob("PPT/*/*.pptx")
        if "svg" not in p.stem.lower() and "优化" not in p.stem and p.parent.name.endswith("1")
    ][0]
    return sample, source


def clear_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public delete API.
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.rId)
        sld_id_lst.remove(sld_id)


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(*(int(hex_color[i : i + 2], 16) for i in (0, 2, 4)))


def inch(value: float):
    return Inches(value)


def add_rect(slide, x, y, w, h, fill, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, inch(x), inch(y), inch(w), inch(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        try:
            shp.line.fill.background()
        except Exception:
            shp.line.color.rgb = fill
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.8)
    return shp


def _set_typeface(r_pr, tag: str, typeface: str):
    child = r_pr.find(qn(f"a:{tag}"))
    if child is None:
        child = OxmlElement(f"a:{tag}")
        r_pr.append(child)
    child.set("typeface", typeface)


def apply_run_typefaces(run, fallback_east_asia: str) -> None:
    r_pr = run._r.get_or_add_rPr()  # noqa: SLF001 - python-pptx has no public mixed-script font API.
    _set_typeface(r_pr, "ea", fallback_east_asia)
    latin = LATIN_FONT if LATIN_RE.search(run.text or "") else fallback_east_asia
    _set_typeface(r_pr, "latin", latin)


def _visual_units(text: str) -> float:
    total = 0.0
    for ch in text:
        if ch.isspace():
            total += 0.35
        elif ch.isascii():
            total += 0.55
        else:
            total += 1.0
    return total


def _splits_latin_token(text: str, index: int) -> bool:
    for match in re.finditer(r"[A-Za-z][A-Za-z0-9_./+\-]*", text):
        if match.start() < index < match.end():
            return True
    return False


def prose_wrap(text: str, limit: float = 20.0) -> str:
    """Wrap mixed Chinese/Latin prose without orphan punctuation or split words."""
    remaining = re.sub(r"\s+", " ", text.replace("\n", "")).strip()
    lines: list[str] = []
    while _visual_units(remaining) > limit:
        valid: list[int] = []
        preferred: list[int] = []
        for index in range(1, len(remaining)):
            if _visual_units(remaining[:index]) > limit:
                break
            if remaining[index] in NO_LINE_START or remaining[index - 1] in NO_LINE_END:
                continue
            if _splits_latin_token(remaining, index):
                continue
            valid.append(index)
            if remaining[index - 1] in PROSE_BREAKS:
                preferred.append(index)
        candidates = preferred or valid or [1]
        balanced = [idx for idx in candidates if _visual_units(remaining[:idx]) >= limit * 0.62]
        cut = (balanced or candidates)[-1]
        lines.append(remaining[:cut].strip())
        remaining = remaining[cut:].strip()
    if remaining:
        lines.append(remaining)
    return "\n".join(lines)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float,
    color=TEXT,
    bold: bool = False,
    font: str = BODY_FONT,
    align=PP_ALIGN.LEFT,
    valign=None,
    margin: float = 0.02,
    fit: bool = False,
    line_spacing: float | None = None,
    wrap: bool = True,
):
    lines = text.split("\n")
    if not fit:
        # Give every text box a little real rendering slack. PowerPoint can
        # visually tolerate tight boxes, but exported slides clip/overlap more
        # easily when the declared box height is only barely one line tall.
        line_factor = 1.34 if font == BODY_FONT else 1.42
        if line_spacing is not None:
            line_factor = max(line_factor, 1.08 * float(line_spacing))
        min_h = max(1, len(lines)) * size * line_factor / 72 + margin * 2 + 0.035
        h = max(h, min_h)
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if fit else MSO_AUTO_SIZE.NONE
    tf.margin_left = inch(margin)
    tf.margin_right = inch(margin)
    tf.margin_top = inch(margin)
    tf.margin_bottom = inch(margin)
    if valign is not None:
        tf.vertical_anchor = valign

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p_pr = p._p.get_or_add_pPr()  # noqa: SLF001 - python-pptx exposes no public mixed-script line-break API.
        p_pr.set("latinLnBrk", "0")
        p_pr.set("eaLnBrk", "1")
        p_pr.set("hangingPunct", "0")
        if line_spacing is not None:
            p.line_spacing = line_spacing
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            apply_run_typefaces(run, font)
    return box


def add_formula(slide, formula: str, x: float, y: float, w: float, h: float, size: float = 17.0, color=GREEN):
    add_rect(slide, x, y, w, h, LIGHT_TAN, None, radius=True)
    return add_text(
        slide,
        formula,
        x + 0.12,
        y + 0.07,
        w - 0.24,
        h - 0.14,
        size,
        color,
        font="Cambria Math",
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
        line_spacing=1.0,
        wrap=False,
    )


def render_latex_formula(name: str, formula: str, color_hex: str = "265B3D") -> bytes:
    formula_dir = OUT_ASSETS / "latex_formulas"
    formula_dir.mkdir(parents=True, exist_ok=True)
    tex_path = formula_dir / f"{name}.tex"
    dvi_path = formula_dir / f"{name}.dvi"
    svg_path = formula_dir / f"{name}.svg"
    tex_source = rf"""\documentclass[border=0pt]{{standalone}}
\usepackage{{amsmath,amssymb}}
\usepackage{{xcolor}}
\begin{{document}}
\color[HTML]{{{color_hex}}}$\displaystyle {formula}$
\end{{document}}
"""
    needs_render = not svg_path.exists() or not tex_path.exists()
    if not needs_render:
        needs_render = tex_path.read_text(encoding="utf-8", errors="ignore") != tex_source
    if needs_render:
        tex_path.write_text(tex_source, encoding="utf-8")
        subprocess.run(
            ["latex", "-interaction=nonstopmode", "-halt-on-error", tex_path.name],
            cwd=formula_dir,
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        subprocess.run(
            ["dvisvgm", "--no-fonts", "--exact", "--bbox=min", f"--output={svg_path.name}", dvi_path.name],
            cwd=formula_dir,
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
    return svg_path.read_bytes()


def add_latex_formula(slide, formula: str, x: float, y: float, w: float, h: float, accent=GREEN, name: str = "formula"):
    add_rect(slide, x, y, w, h, LIGHT_TAN, None, radius=True)
    color_hex = f"{FORMULA_GREEN[0]:02X}{FORMULA_GREEN[1]:02X}{FORMULA_GREEN[2]:02X}"
    formula_png = trim_whitespace_bytes(
        normalize_image_bytes(render_latex_formula(name, formula, color_hex)),
        tolerance=250,
        pad=4,
    )
    im = Image.open(io.BytesIO(formula_png))
    max_w = w - 0.36
    target_heights = {
        "normalize": 0.46,
        "relu": 0.31,
        "maxpool": 0.41,
        "softmax": 0.52,
        "cross_entropy": 0.50,
        "convolution": 0.56,
        "accuracy": 0.78,
        "adamw": 0.74,
    }
    target_h = min(h - 0.04, target_heights.get(name, 0.56))
    target_w = target_h * im.width / im.height
    if target_w > max_w:
        target_w = max_w
        target_h = target_w * im.height / im.width
    left = x + (w - target_w) / 2
    top = y + (h - target_h) / 2
    return slide.shapes.add_picture(io.BytesIO(formula_png), inch(left), inch(top), width=inch(target_w), height=inch(target_h))


def add_label(slide, text, x, y, w, h, color=GREEN, size=10):
    shp = add_rect(slide, x, y, w, h, BG, color, radius=True)
    add_text(
        slide,
        text,
        x + 0.05,
        y + 0.015,
        w - 0.1,
        h - 0.03,
        size,
        color=color,
        font=MONO_FONT,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )
    return shp


def add_bg(slide, deco: str = "default"):
    add_rect(slide, 0, 0, 13.333, 7.5, BG, None, radius=False)
    assets = prepare_decor_assets()
    if deco == "cover":
        add_picture_fit(slide, assets["leaf_large"], 9.15, -0.15, 3.85, 4.7)
        add_picture_fit(slide, assets["leaf_line"], 0.85, 4.25, 1.85, 2.1)
    elif deco == "close":
        add_picture_fit(slide, assets["leaf_large"], 2.25, -1.05, 8.9, 9.1)
    else:
        add_picture_fit(slide, assets["leaf_line"], 11.2, -0.15, 1.65, 1.35)
        add_picture_fit(slide, assets["leaf_line"], -0.22, 5.75, 1.45, 1.2)


def add_header(slide, eyebrow: str, title: str, subtitle: str, page: str | None = None):
    add_rect(slide, 0.42, 0.44, 0.08, 0.62, GREEN, None, radius=True)
    add_text(slide, eyebrow, 0.62, 0.46, 4.8, 0.22, 10.5, MUTED, font=MONO_FONT, fit=False)
    add_text(slide, title, 0.62, 0.78, 8.4, 0.54, 25, GREEN, font=TITLE_FONT)
    if subtitle:
        add_text(slide, subtitle, 0.63, 1.34, 8.8, 0.24, 11.5, MUTED, font=BODY_FONT)
    if page:
        add_text(slide, page, 11.55, 6.95, 1.1, 0.25, 9.5, MUTED, font=MONO_FONT, align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title=None, accent=GREEN, fill=CARD, pale=False):
    add_rect(slide, x, y, w, h, CARD_2 if pale else fill, LINE, radius=True)
    add_rect(slide, x, y, w, 0.045, accent, accent, radius=False)
    if title:
        add_text(slide, title, x + 0.22, y + 0.22, w - 0.44, 0.34, 16.5, accent, font=TITLE_FONT, bold=True)


def add_inner(slide, x, y, w, h, text, size=13.2, title=None, accent=GREEN, fill=CARD_2, text_margin=0.10):
    add_rect(slide, x, y, w, h, fill, None, radius=True)
    inner_margin = text_margin
    if title:
        add_text(slide, title, x + inner_margin, y + 0.12, w - inner_margin * 2, 0.26, size + 0.5, accent, bold=True)
        body_y = y + 0.46
        body_h = h - 0.56
    else:
        body_y = y + 0.12
        body_h = h - 0.22
    add_text(slide, text, x + inner_margin, body_y, w - inner_margin * 2, body_h, size, TEXT, font=BODY_FONT, line_spacing=1.18)


def add_picture_fit(slide, source, x, y, w, h, pad=0.0, trim=False):
    if isinstance(source, (str, Path)):
        img_bytes = Path(source).read_bytes()
    else:
        img_bytes = source
    img_bytes = normalize_image_bytes(img_bytes)
    if trim:
        img_bytes = trim_whitespace_bytes(img_bytes)
    im = Image.open(io.BytesIO(img_bytes))
    box_w, box_h = w - 2 * pad, h - 2 * pad
    img_ratio = im.width / im.height
    box_ratio = box_w / box_h
    if img_ratio > box_ratio:
        new_w = box_w
        new_h = box_w / img_ratio
    else:
        new_h = box_h
        new_w = box_h * img_ratio
    left = x + pad + (box_w - new_w) / 2
    top = y + pad + (box_h - new_h) / 2
    return slide.shapes.add_picture(io.BytesIO(img_bytes), inch(left), inch(top), width=inch(new_w), height=inch(new_h))


def normalize_image_bytes(img_bytes: bytes) -> bytes:
    stripped = img_bytes.lstrip()
    if stripped.startswith(b"<svg") or (stripped.startswith(b"<?xml") and b"<svg" in stripped[:300]):
        drawing = svg2rlg(io.BytesIO(img_bytes))
        drawing.scale(4, 4)
        drawing.width *= 4
        drawing.height *= 4
        out = io.BytesIO()
        renderPM.drawToFile(drawing, out, fmt="PNG")
        return white_to_transparent(out.getvalue())
    return img_bytes


def white_to_transparent(img_bytes: bytes, threshold: int = 248) -> bytes:
    im = Image.open(io.BytesIO(img_bytes)).convert("RGBA")
    data = []
    for r, g, b, a in im.getdata():
        if r >= threshold and g >= threshold and b >= threshold:
            data.append((255, 255, 255, 0))
        else:
            data.append((r, g, b, a))
    im.putdata(data)
    out = io.BytesIO()
    im.save(out, format="PNG")
    return out.getvalue()


def trim_whitespace_bytes(img_bytes: bytes, tolerance: int = 246, pad: int = 18) -> bytes:
    im = Image.open(io.BytesIO(img_bytes)).convert("RGBA")
    pixels = im.load()
    min_x, min_y = im.width, im.height
    max_x, max_y = -1, -1
    for y in range(im.height):
        for x in range(im.width):
            r, g, b, a = pixels[x, y]
            if a > 16 and min(r, g, b) < tolerance:
                min_x = min(min_x, x)
                min_y = min(min_y, y)
                max_x = max(max_x, x)
                max_y = max(max_y, y)
    if max_x < min_x or max_y < min_y:
        return img_bytes
    min_x = max(0, min_x - pad)
    min_y = max(0, min_y - pad)
    max_x = min(im.width - 1, max_x + pad)
    max_y = min(im.height - 1, max_y + pad)
    cropped = im.crop((min_x, min_y, max_x + 1, max_y + 1))
    out = io.BytesIO()
    cropped.save(out, format="PNG")
    return out.getvalue()


def prepare_decor_assets() -> dict[str, bytes]:
    OUT_ASSETS.mkdir(parents=True, exist_ok=True)
    large = OUT_ASSETS / "leaf_large_faded.png"
    line = OUT_ASSETS / "leaf_line_faded.png"
    if large.exists() and line.exists():
        return {"leaf_large": large.read_bytes(), "leaf_line": line.read_bytes()}

    sample, _ = locate_source()
    with zipfile.ZipFile(sample) as zf:
        large_src = Image.open(io.BytesIO(zf.read("ppt/media/image-1-1.jpg"))).convert("RGB")
        line_src = Image.open(io.BytesIO(zf.read("ppt/media/image-2-1.jpg"))).convert("RGB")

    def fade(im: Image.Image, alpha: float) -> Image.Image:
        bg = Image.new("RGB", im.size, (248, 247, 241))
        return Image.blend(bg, im, alpha)

    fade(large_src, 0.18).save(large)
    fade(line_src, 0.12).save(line)
    return {"leaf_large": large.read_bytes(), "leaf_line": line.read_bytes()}


def walk_pictures(shapes):
    for shp in shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk_pictures(shp.shapes)
        elif shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shp


def picture_blob(pic) -> bytes | None:
    try:
        return pic.image.blob
    except ValueError:
        blip = pic._element.blipFill.blip  # noqa: SLF001 - fallback for linked/legacy picture nodes.
        r_id = getattr(blip, "rEmbed", None) or getattr(blip, "rLink", None)
        if not r_id:
            match = re.search(r'r:embed="([^"]+)"', blip.xml)
            r_id = match.group(1) if match else None
        if not r_id:
            return None
        try:
            return pic.part.related_part(r_id).blob
        except Exception:
            return None


def extract_images(source_path: Path) -> dict[int, list[bytes]]:
    prs = Presentation(str(source_path))
    images: dict[int, list[bytes]] = {}
    for i, slide in enumerate(prs.slides, 1):
        images[i] = [blob for pic in walk_pictures(slide.shapes) if (blob := picture_blob(pic))]
    return images


def add_footer_ref(slide, ref: str, page: str):
    if ref:
        add_text(slide, ref, 0.7, 6.92, 9.6, 0.22, 6.5, RGBColor(160, 158, 149), font=BODY_FONT)
    add_text(slide, page, 11.55, 6.92, 1.1, 0.24, 8.8, MUTED, font=MONO_FONT, align=PP_ALIGN.RIGHT)


def build_deck() -> Path:
    sample_path, source_path = locate_source()
    algorithm_flow = figure_path("algorithm_flow_academic.png")
    training_curve = figure_path("training_curve_20260426_vector.svg")
    confusion_matrix = figure_path("confusion_matrix_20260426_vector.svg")

    prs = Presentation(str(sample_path))
    clear_slides(prs)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[0]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    add_bg(slide, "cover")
    add_label(slide, "DEEP LEARNING COURSE EXPERIMENT", 0.42, 1.15, 3.62, 0.42, GREEN, 9.5)
    add_text(slide, "实验一：手写数字字符识别", 0.45, 2.2, 7.9, 0.95, 39, GREEN, font=TITLE_FONT, fit=True, wrap=False)
    add_text(slide, "基于PyTorch卷积神经网络实现MNIST数据集分类", 0.47, 3.62, 6.8, 0.38, 16, TEXT)
    add_rect(slide, 0.42, 4.43, 1.0, 0.025, GREEN, None, radius=False)
    add_text(slide, "深度学习课程实验", 1.62, 4.32, 2.6, 0.28, 15, TEXT, font=TITLE_FONT)
    add_text(slide, "2026年04月", 1.62, 4.68, 1.8, 0.24, 10.5, MUTED, font=MONO_FONT)
    add_text(slide, "卷积神经网络 · MNIST数据集 · PyTorch框架 · 准确率99.28%", 0.42, 6.15, 6.4, 0.24, 9.5, MUTED)
    add_rect(slide, 0.42, 6.43, 12.0, 0.01, LINE, None, radius=False)
    for idx, label in enumerate(["课程实验", "模型实现", "结果分析", "总结展望"]):
        x = 0.47 + idx * 2.25
        add_rect(slide, x, 6.75, 0.16, 0.16, GREEN if idx in (0, 3) else TAN, None, radius=True)
        add_text(slide, label, x + 0.25, 6.68, 1.0, 0.25, 8.5, MUTED)
    add_text(slide, "01 / 11", 11.55, 6.78, 1.1, 0.25, 8.8, MUTED, font=MONO_FONT, align=PP_ALIGN.RIGHT)

    # 2. Overview
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "EXPERIMENT OVERVIEW", "实验概述", "目标要求 · 数据集 · 评估标准 · 完成情况")
    cards = [
        ("运行环境", "任务要求\n配置 PyTorch 框架\n支持 GPU 加速训练\n\n完成状态\nPython 3.12.12\nPyTorch 2.9.0+cu128\nCUDA 12.8 正常"),
        ("模型架构", "任务要求\n设计卷积神经网络\n标准层级结构\n\n完成状态\n两层卷积 + ReLU\n+ 最大池化层\n+ 全连接输出"),
        ("数据集", "任务要求\n使用 MNIST 手写\n数字图像数据库\n\n完成状态\n60000 训练样本\n10000 测试样本\n28x28 像素, 10分类"),
        ("性能指标", "任务要求\n测试集准确率 >= 98%\n\n完成状态\n最佳测试准确率\n99.28%"),
    ]
    for i, (title, body) in enumerate(cards):
        x = 0.42 + i * 3.16
        add_card(slide, x, 1.78, 2.9, 2.65, title, [GREEN, SAGE, TAN, GREEN][i])
        add_text(slide, f"{i+1:02d}", x + 2.22, 2.02, 0.5, 0.34, 14, [GREEN, SAGE, TAN, GREEN][i], font=MONO_FONT, align=PP_ALIGN.RIGHT)
        add_text(slide, body, x + 0.23, 2.32, 2.46, 1.84, 10.5, TEXT, line_spacing=0.98)
        add_rect(slide, x + 2.02, 4.02, 0.68, 0.27, CARD_2, None, radius=True)
        add_text(slide, "完成", x + 2.08, 4.055, 0.56, 0.2, 7.8, GREEN, bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_card(slide, 0.42, 4.72, 12.47, 1.82, "实验目的与数据说明", GREEN, fill=CARD_2, pale=True)
    add_text(
        slide,
        "实验目的：深入掌握卷积神经网络核心原理，利用PyTorch构建网络模型，在MNIST数据集上完成训练验证，准确率需达到98%以上。\n\nMNIST数据集简介：手写数字识别领域广泛使用的公开基准数据。包含60000幅训练图像和10000幅测试图像，已完成尺寸标准化和中心化预处理。",
        0.68,
        5.18,
        11.9,
        1.02,
        12.0,
        TEXT,
        line_spacing=1.08,
    )
    add_footer_ref(slide, "[1] LeCun et al., The MNIST database of handwritten digits  [2] Gradient-based learning applied to document recognition, 1998", "02 / 11")

    # 3-4. Theory
    theory_pages = [
        (
            "理论基础（一）",
            "输入标准化 · 卷积特征 · 非线性激活 · 空间池化",
            "03 / 11",
            "[1] LeCun et al., The MNIST database of handwritten digits  [3] Nair & Hinton, Rectified linear units improve RBMs, ICML 2010",
            [
                ("normalize", 0.55, 1.76, "数据归一化", "按通道减均值并除以标准差，\n减小输入尺度差异，使训练更稳定。", r"\tilde{x} = \frac{x-\mu}{\sigma}", GREEN),
                ("convolution", 6.75, 1.76, "卷积运算", "共享卷积核在局部邻域滑动计算，\n用较少参数提取笔画边缘与结构特征。", r"\begin{aligned}S(i,j)&=(I * K)(i,j)\\&=\sum_m\sum_n I(i+m,j+n)K(m,n)\end{aligned}", SAGE),
                ("relu", 0.55, 4.18, "ReLU 激活", "保留正向响应并抑制负响应，引入非线性，增强网络表达能力。", r"\mathrm{ReLU}(z)=\max(0,z)", TAN),
                ("maxpool", 6.75, 4.18, "最大池化", "在局部窗口保留最大激活，\n降低空间分辨率并增强轻微平移鲁棒性。", r"y_{i,j} = \max_{(m,n)\in\Omega_{i,j}} x_{m,n}", GREEN),
            ],
        ),
        (
            "理论基础（二）",
            "概率建模 · 损失函数 · 评估指标 · 参数优化",
            "04 / 11",
            "[3] Paszke et al., PyTorch: An imperative style, high-performance DL library, NeurIPS 2019  [5] Loshchilov & Hutter, Decoupled weight decay regularization, ICLR 2019",
            [
                ("softmax", 0.55, 1.76, "Softmax 概率", "将 logits 转换为归一化概率，\n便于解释模型对不同数字的类别置信度。", r"p_{i,c} = \frac{\exp(z_{i,c})}{\sum_{k=1}^{10}\exp(z_{i,k})}", SAGE),
                ("cross_entropy", 6.75, 1.76, "交叉熵损失", "对真实类别概率取负对数，\n推动模型提高正确类别的预测概率。", r"\mathcal{L}(\theta) = -\frac{1}{N}\sum_{i=1}^{N}\log p_{i,y_i}", TAN),
                ("accuracy", 0.55, 4.18, "准确率指标", "先取 logits 最大的类别作为预测，\n再统计与真实标签一致的样本比例。", r"\begin{aligned}\hat{y}_i&=\arg\max_c f_\theta(x_i)_c\\\mathrm{Acc}&=\frac{1}{N}\sum_{i=1}^{N}\mathbb{I}(\hat{y}_i=y_i)\end{aligned}", GREEN),
                ("adamw", 6.75, 4.18, "AdamW 优化", "将权重衰减与自适应梯度更新解耦，\n控制模型复杂度并改善泛化。", r"\begin{aligned}\theta_{t+1}&=\theta_t-\eta\frac{\hat{m}_t}{\sqrt{\hat{v}_t}+\epsilon}\\[-0.4em]&\quad-\eta\lambda\theta_t\end{aligned}", SAGE),
            ],
        ),
    ]
    for title, subtitle, page, ref, theory_cards in theory_pages:
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_header(slide, "THEORETICAL BASIS", title, subtitle)
        for name, x, y, card_title, desc, formula_text, accent in theory_cards:
            add_rect(slide, x, y, 5.75, 2.14, CARD, LINE, radius=True)
            add_rect(slide, x, y, 5.75, 0.045, accent, accent, radius=False)
            add_text(slide, card_title, x + 0.28, y + 0.20, 2.4, 0.28, 16.2, accent, font=TITLE_FONT, bold=True)
            add_text(slide, desc, x + 0.28, y + 0.60, 5.18, 0.56, 10.6, MUTED, line_spacing=1.0)
            add_latex_formula(slide, formula_text, x + 0.34, y + 1.24, 5.07, 0.76, accent, name)
        add_footer_ref(slide, ref, page)

    # 5. Architecture
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "MODEL PIPELINE", "网络结构与训练流程", "MNIST CNN · 特征提取、分类头与反向传播闭环", "05 / 11")
    add_card(slide, 0.55, 1.78, 12.25, 4.95, None, GREEN)
    add_picture_fit(slide, algorithm_flow, 0.72, 2.0, 11.92, 4.50, pad=0.02, trim=True)

    # 6. Setup
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "EXPERIMENT SETUP", "实验环境与超参数配置", "软硬件环境 · 训练超参数 · 优化器与损失函数")
    add_card(slide, 0.55, 1.78, 5.85, 3.85, "运行环境", GREEN)
    env_rows = [("Python版本", "3.12.12"), ("PyTorch版本", "2.9.0+cu128"), ("torchvision版本", "0.24.0+cu128"), ("CUDA版本", "12.8"), ("CUDA可用", "True"), ("cuDNN版本", "91002")]
    for i, (k, v) in enumerate(env_rows):
        y = 2.42 + i * 0.46
        add_text(slide, k, 0.9, y, 2.2, 0.24, 10.2, MUTED)
        add_text(slide, v, 4.55, y, 1.45, 0.24, 10.8, GREEN if i in (3, 4) else TEXT, bold=i in (3, 4), align=PP_ALIGN.RIGHT)
        add_rect(slide, 0.9, y + 0.34, 4.9, 0.006, LINE, None, radius=False)
    add_card(slide, 6.75, 1.78, 5.75, 3.85, "训练超参数", SAGE)
    hp_rows = [("Epochs", "5"), ("Batch Size", "128"), ("Learning Rate", "0.001"), ("Weight Decay", "1e-4"), ("Dropout", "0.1"), ("Random Seed", "42")]
    for i, (k, v) in enumerate(hp_rows):
        y = 2.42 + i * 0.46
        add_text(slide, k, 7.08, y, 2.3, 0.24, 10.0, MUTED, font=MONO_FONT)
        add_text(slide, v, 11.05, y, 0.95, 0.24, 11.0, GREEN, bold=True, align=PP_ALIGN.RIGHT)
        add_rect(slide, 7.08, y + 0.34, 4.55, 0.006, LINE, None, radius=False)
    add_inner(slide, 0.55, 5.82, 5.85, 0.86, "权重衰减与梯度更新分离，提升泛化效果 (Loshchilov & Hutter, ICLR 2019)", 8.7, "优化器: AdamW", GREEN)
    add_inner(slide, 6.75, 5.82, 5.75, 0.86, "输入原始logits和目标类别索引 (PyTorch NeurIPS 2019)", 8.7, "损失函数: 交叉熵损失", TAN)
    add_footer_ref(slide, "[3] Paszke et al., PyTorch: An imperative style, high-performance DL library, NeurIPS 2019  [5] Loshchilov & Hutter, Decoupled weight decay, ICLR 2019", "06 / 11")

    # 7. Results
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "TRAINING RESULTS", "训练结果与验证数据", "5轮训练 · 最优测试准确率99.28% · 超过98%目标", "07 / 11")
    add_card(slide, 0.55, 1.78, 5.1, 3.78, "训练过程概览", GREEN)
    headers = ["轮次", "训练损失", "训练准确率", "测试准确率"]
    xs = [0.85, 1.85, 3.05, 4.25]
    for x, h in zip(xs, headers):
        add_text(slide, h, x, 2.36, 0.9, 0.22, 8.8, GREEN, bold=True, align=PP_ALIGN.CENTER)
    rows = [
        ("1", "0.2726", "91.63%", "98.31%"),
        ("2", "0.0911", "97.23%", "98.78%"),
        ("3", "0.0702", "97.86%", "98.95%"),
        ("4", "0.0587", "98.22%", "99.12%"),
        ("5", "0.0537", "98.37%", "99.28%"),
    ]
    for i, row in enumerate(rows):
        y = 2.78 + i * 0.36
        if i == 4:
            add_rect(slide, 0.78, y - 0.04, 4.55, 0.32, CARD_2, None, radius=True)
        for x, val in zip(xs, row):
            add_text(slide, val, x, y, 0.9, 0.22, 9.1, GREEN if val.endswith("%") and val.startswith("99") else TEXT, bold=i == 4, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.78, 4.72, 4.55, 0.006, GREEN, None, radius=False)
    add_text(slide, "核心发现", 0.85, 4.84, 1.2, 0.24, 11.0, GREEN, bold=True)
    add_text(slide, "• 第1轮即突破98%门槛 (98.31%)\n• 训练损失: 0.2726 → 0.0537 (↓80%)", 0.85, 5.13, 4.3, 0.38, 9.0, TEXT)
    add_card(slide, 5.92, 1.78, 6.6, 4.68, "训练验证曲线", SAGE)
    add_picture_fit(slide, training_curve, 6.18, 2.2, 6.1, 3.70, pad=0.05, trim=True)

    # 8. Analysis
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "RESULT ANALYSIS", "结果分析与达标判定", "混淆矩阵分析 · 目标对比 · 错误模式", "08 / 11")
    add_card(slide, 0.55, 1.78, 5.55, 4.55, "测试集混淆矩阵", GREEN)
    add_picture_fit(slide, confusion_matrix, 0.78, 2.2, 5.1, 3.8, pad=0.04, trim=True)
    add_card(slide, 6.35, 1.78, 6.15, 4.55, "目标与实际", SAGE)
    add_text(slide, "目标准确率", 7.05, 2.3, 1.4, 0.22, 9.8, MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "98%", 7.05, 2.56, 1.4, 0.6, 25, MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "→", 8.75, 2.63, 0.5, 0.44, 16, GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, "实际准确率", 9.65, 2.3, 1.6, 0.22, 9.8, MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "99.28%", 9.55, 2.54, 1.9, 0.62, 25, GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 11.16, 3.1, 0.82, 0.24, CARD_2, None, radius=True)
    add_text(slide, "+1.28%", 11.18, 3.12, 0.78, 0.16, 8.5, GREEN, bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_rect(slide, 6.72, 3.54, 5.45, 0.006, LINE, None, radius=False)
    add_text(slide, "错误模式研究", 6.75, 3.68, 1.8, 0.28, 15, TAN, font=TITLE_FONT, bold=True)
    add_text(slide, "大部分样本位于主对角线，模型对各数字类别具备良好分辨力。", 6.75, 4.03, 5.3, 0.25, 9.8, TEXT)
    for i, (main, sub) in enumerate([("3 ↔ 5", "形态相似"), ("4 ↔ 9", "顶部闭合"), ("7 ↔ 9", "笔画相近")]):
        x = 6.8 + i * 1.78
        add_rect(slide, x, 4.48, 1.48, 0.48, LIGHT_TAN, None, radius=True)
        add_text(slide, main, x, 4.52, 1.48, 0.24, 8.6, ORANGE, bold=True, align=PP_ALIGN.CENTER, margin=0)
        add_text(slide, sub, x, 4.75, 1.48, 0.16, 7.5, MUTED, align=PP_ALIGN.CENTER, margin=0)
    add_text(slide, "与手写数字识别研究的经典误差分布一致。", 6.75, 5.14, 5.1, 0.22, 8.8, MUTED)
    add_text(slide, "各类别准确率：数字0: 99.9% │ 数字1: 99.8% │ 数字2: 99.3% │ 数字3: 98.3% │ 数字4: 99.4%\n数字5: 99.4% │ 数字6: 99.1% │ 数字7: 99.7% │ 数字8: 98.6% │ 数字9: 98.8%", 6.75, 5.47, 5.25, 0.4, 7.6, TEXT)

    # 9. Discussion
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "DISCUSSION & OUTLOOK", "讨论与展望", "从实验结论走向局限分析、应用外延与后续研究", "09 / 11")
    discussion = [
        ("局限性", "MNIST 背景干净、尺寸统一，\n无法覆盖真实票据、表格扫描、\n拍照场景中的噪声、倾斜、遮挡\n与分布漂移。", GREEN),
        ("应用方向", "可扩展到答题卡识别、\n票据数字读取、仪表盘编号识别\n和轻量 OCR 系统中的\n局部字符分类模块。", SAGE),
        ("深入研究", "后续应加入鲁棒数据增强、\n置信度校准、拒识机制、\n错例聚类和轻量化部署评估，\n而不仅追求更高测试准确率。", TAN),
    ]
    for i, (title, body, accent) in enumerate(discussion):
        x = 0.45 + i * 4.25
        add_card(slide, x, 1.78, 3.95, 5.18, None, accent)
        add_rect(slide, x + 0.28, 2.1, 0.66, 0.66, CARD_2, LINE, radius=True)
        add_text(slide, f"{i+1}", x + 0.28, 2.24, 0.66, 0.32, 14.0, accent, font=MONO_FONT, bold=True, align=PP_ALIGN.CENTER, margin=0)
        add_text(slide, title, x + 1.08, 2.16, 2.25, 0.36, 18.0, accent, font=TITLE_FONT, bold=True)
        add_inner(slide, x + 0.09, 3.02, 3.77, 2.92, prose_wrap(body, 20.0), 12.3, accent=accent, fill=CARD_2, text_margin=0.08)

    # 10. Summary
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "SUMMARY", "总结", "实验目标、结果证据与核心收获", "10 / 11")
    summary_items = [
        ("1", "完成 PyTorch/CUDA\n环境验证、CNN 模型训练、\ncheckpoint 复评\n和可视化分析。"),
        ("2", "测试准确率达到 99.28%，\n显著超过 98% 的\n课程验收目标。"),
        ("3", "实验重点在于建立可复现的\n图像分类工程链路，\n而不仅是获得单次高分。"),
    ]
    for i, (num, body) in enumerate(summary_items):
        x = 0.45 + i * 4.25
        add_card(slide, x, 1.78, 3.95, 5.18, None, [GREEN, SAGE, TAN][i])
        add_rect(slide, x + 0.28, 2.1, 0.66, 0.66, CARD_2, LINE, radius=True)
        add_text(slide, num, x + 0.28, 2.24, 0.66, 0.32, 14.0, [GREEN, SAGE, TAN][i], font=MONO_FONT, bold=True, align=PP_ALIGN.CENTER, margin=0)
        add_text(slide, ["工程闭环", "结果证据", "核心收获"][i], x + 1.08, 2.16, 2.0, 0.36, 18.0, [GREEN, SAGE, TAN][i], font=TITLE_FONT, bold=True)
        add_inner(slide, x + 0.09, 3.02, 3.77, 2.92, prose_wrap(body, 20.0), 12.8, accent=[GREEN, SAGE, TAN][i], fill=CARD_2, text_margin=0.08)

    # 11. Closing
    slide = prs.slides.add_slide(blank)
    add_bg(slide, "close")
    add_rect(slide, 6.23, 0.62, 0.86, 0.86, CARD_2, LINE, radius=True)
    add_text(slide, "CNN", 6.25, 0.88, 0.82, 0.32, 10.5, GREEN, font=MONO_FONT, bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_text(slide, "实验一：手写数字字符识别", 2.7, 1.95, 7.9, 0.65, 31, GREEN, font=TITLE_FONT, align=PP_ALIGN.CENTER)
    add_text(slide, "感谢聆听！", 4.1, 2.74, 5.2, 0.78, 30, GREEN, font=TITLE_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 3.75, 3.72, 1.35, 0.02, GREEN, None, radius=False)
    add_text(slide, "敬请老师和同学批评指正", 5.3, 3.56, 2.8, 0.36, 15, TEXT, font=TITLE_FONT, align=PP_ALIGN.CENTER)
    add_rect(slide, 8.25, 3.72, 1.35, 0.02, GREEN, None, radius=False)
    add_rect(slide, 3.65, 4.6, 6.1, 0.74, CARD, LINE, radius=True)
    add_text(slide, "Deep Learning Course Experiment Presentation", 3.9, 4.86, 5.6, 0.23, 11, MUTED, font=MONO_FONT, align=PP_ALIGN.CENTER)
    add_text(slide, "11 / 11", 11.55, 6.82, 1.1, 0.24, 8.8, MUTED, font=MONO_FONT, align=PP_ALIGN.RIGHT)

    out = source_path.with_name(f"{source_path.stem}_优化版.pptx")
    fallback_out = source_path.with_name(f"{source_path.stem}_优化版_细节修正版.pptx")
    prs.core_properties.title = "实验一：手写数字字符识别（优化版）"
    try:
        prs.save(out)
        return out
    except PermissionError:
        prs.save(fallback_out)
        return fallback_out


if __name__ == "__main__":
    print(build_deck())
