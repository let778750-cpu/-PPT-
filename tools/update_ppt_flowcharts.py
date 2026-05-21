from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]


EXPERIMENTS = {
    "exp1": {
        "markers": ("实验一", "实验1", "experiment1_mnist"),
        "title": "网络结构与训练流程",
        "subtitle": "MNIST CNN · 数据加载、优化、评估与可复现产物",
        "image": "实验报告/实验1/figures/algorithm_flow_20260426_vector.png",
    },
    "exp2": {
        "markers": ("实验2", "experiment2_vit"),
        "title": "算法设计与训练流程",
        "subtitle": "CIFAR-10 Vision Transformer · Patch tokenization 到 checkpoint",
        "image": "实验报告/实验2/figures/algorithm_flow_academic.png",
    },
    "exp3": {
        "markers": ("实验3", "自动写诗"),
        "title": "算法设计与生成流程",
        "subtitle": "LSTM 自动写诗 · 语料处理、训练闭环与温度采样",
        "image": "实验报告/实验3/figures/algorithm_flow_academic.png",
    },
    "exp4": {
        "markers": ("实验4", "experiment4_nmt"),
        "title": "算法设计与翻译流程",
        "subtitle": "Transformer NMT · 编码解码、Beam Search 与 BLEU4 评估",
        "image": "实验报告/实验4/figures/algorithm_flow_academic.png",
    },
    "exp6": {
        "markers": ("实验6", "experiment6_segnet"),
        "title": "算法设计与分割流程",
        "subtitle": "SegNet 街景分割 · 数据准备、像素级评估与样例预测",
        "image": "实验报告/实验6/figures/algorithm_flow_academic.png",
    },
    "exp7": {
        "markers": ("实验7", "experiment7_lstm"),
        "title": "算法设计与语言建模流程",
        "subtitle": "PTB LSTM LM · BPTT、验证集选择与测试 PPL",
        "image": "实验报告/实验7/figures/algorithm_flow.png",
    },
}


SEARCH_ROOTS = [
    "PPT",
    "实验报告和PPT（新）/PPT",
    "PPT生成skill/projects",
    "projects",
    "code/work7 code/projects",
]


def classify(path: Path) -> dict[str, str] | None:
    text = str(path)
    for spec in EXPERIMENTS.values():
        if any(marker in text for marker in spec["markers"]):
            return spec
    return None


def collect_targets() -> list[tuple[Path, dict[str, str]]]:
    targets: list[tuple[Path, dict[str, str]]] = []
    seen: set[Path] = set()
    for rel_root in SEARCH_ROOTS:
        root = ROOT / rel_root
        if not root.exists():
            continue
        for path in root.rglob("*.pptx"):
            spec = classify(path)
            if spec is None:
                continue
            resolved = path.resolve()
            if resolved in seen:
                continue
            targets.append((path, spec))
            seen.add(resolved)
    return sorted(targets, key=lambda item: str(item[0]))


def slide_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            parts.append(shape.text)
    return "\n".join(parts)


def choose_slide(prs: Presentation) -> int:
    preferred = ("算法设计", "网络结构与训练流程", "训练流程", "执行流程")
    for idx, slide in enumerate(prs.slides):
        if idx == 0:
            continue
        text = slide_text(slide)
        if any(token in text for token in preferred):
            return idx
    return min(3, len(prs.slides) - 1)


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


def add_textbox(slide, left, top, width, height, text, font_size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def update_deck(path: Path, spec: dict[str, str]) -> tuple[int, int]:
    prs = Presentation(path)
    idx = choose_slide(prs)
    slide = prs.slides[idx]
    clear_slide(slide)

    sw = prs.slide_width
    sh = prs.slide_height

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(247, 249, 252)
    bg.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.44), Inches(0.09), Inches(0.58))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(37, 99, 235)
    accent.line.fill.background()

    add_textbox(
        slide,
        Inches(0.75),
        Inches(0.30),
        sw - Inches(1.5),
        Inches(0.36),
        spec["title"],
        24,
        (15, 23, 42),
        bold=True,
    )
    add_textbox(
        slide,
        Inches(0.75),
        Inches(0.72),
        sw - Inches(1.5),
        Inches(0.30),
        spec["subtitle"],
        10.5,
        (100, 116, 139),
    )

    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.18), sw - Inches(1.1), Pt(1.2))
    rule.fill.solid()
    rule.fill.fore_color.rgb = RGBColor(203, 213, 225)
    rule.line.fill.background()

    image_path = ROOT / spec["image"]
    pic_top = Inches(1.42)
    max_width = sw - Inches(1.1)
    max_height = sh - pic_top - Inches(0.42)
    with Image.open(image_path) as img:
        aspect = img.width / img.height
    pic_width = max_width
    pic_height = int(pic_width / aspect)
    if pic_height > max_height:
        pic_height = max_height
        pic_width = int(pic_height * aspect)
    pic_left = int((sw - pic_width) / 2)
    slide.shapes.add_picture(str(image_path), pic_left, pic_top, width=pic_width, height=pic_height)

    prs.save(path)
    return idx + 1, len(prs.slides)


def main() -> None:
    for path, spec in collect_targets():
        slide_no, slide_count = update_deck(path, spec)
        print(f"updated slide {slide_no}/{slide_count}: {path.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
